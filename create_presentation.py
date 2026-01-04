#!/usr/bin/env python3
"""
Palli Sahayak Presentation Generator
Creates a professional PPTX presentation about the Voice AI Helpline
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme as hex tuples (R, G, B) - Professional palliative care theme
COLORS = {
    'primary': (0x2E, 0x7D, 0x32),      # Green (compassion, healing)
    'secondary': (0x00, 0x69, 0x6B),    # Teal
    'accent': (0xFF, 0x8F, 0x00),       # Amber
    'dark': (0x21, 0x21, 0x21),         # Dark gray
    'light': (0xF5, 0xF5, 0xF5),        # Light gray
    'white': (0xFF, 0xFF, 0xFF),
    'red': (0xC6, 0x28, 0x28),          # For crisis stats
    'blue': (0x15, 0x65, 0xC0),         # For technology
}

def set_shape_color(shape, color_tuple):
    """Set the fill color of a shape."""
    from pptx.dml.color import RGBColor
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color_tuple)

def set_text_color(run_or_paragraph, color_tuple):
    """Set the color of text."""
    from pptx.dml.color import RGBColor
    run_or_paragraph.font.color.rgb = RGBColor(*color_tuple)

def set_line_color(shape, color_tuple):
    """Set the line/border color of a shape."""
    from pptx.dml.color import RGBColor
    shape.line.color.rgb = RGBColor(*color_tuple)

def create_presentation():
    """Create the main presentation."""
    prs = Presentation()
    prs.slide_width = Inches(16)  # Widescreen 16:9
    prs.slide_height = Inches(9)

    # Add slides
    add_title_slide(prs)
    add_acknowledgements_slide(prs)
    add_crisis_slide(prs)
    add_opportunity_slide(prs)
    add_introducing_slide(prs)
    add_architecture_slide(prs)
    add_voice_first_slide(prs)
    add_languages_slide(prs)
    add_rag_foundation_slide(prs)
    add_knowledge_architecture_slide(prs)
    add_clinical_validation_slide(prs)
    add_personalization_slide(prs)
    add_analytics_slide(prs)
    add_dpg_slide(prs)
    add_global_replicability_slide(prs)
    add_impact_metrics_slide(prs)
    add_partners_slide(prs)
    add_tech_stack_slide(prs)
    add_indian_context_slide(prs)
    add_roadmap_slide(prs)
    add_call_to_action_slide(prs)
    add_thank_you_slide(prs)

    return prs

def add_title_text(slide, text, top=Inches(0.5), font_size=44, color=None):
    """Add title text to a slide."""
    title = slide.shapes.add_textbox(Inches(0.5), top, Inches(15), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    if color:
        set_text_color(p, color)
    return title

def add_subtitle_text(slide, text, top=Inches(1.5), font_size=24, color=None):
    """Add subtitle text to a slide."""
    subtitle = slide.shapes.add_textbox(Inches(0.5), top, Inches(15), Inches(0.8))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.alignment = PP_ALIGN.CENTER
    if color:
        set_text_color(p, color)
    return subtitle

def add_title_slide(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
    set_shape_color(background, COLORS['primary'])
    background.line.fill.background()

    # Main title
    title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Palli Sahayak"
    p.font.size = Pt(72)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    # Hindi name
    hindi = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(14), Inches(0.6))
    tf = hindi.text_frame
    p = tf.paragraphs[0]
    p.text = "पल्ली सहायक"
    p.font.size = Pt(36)
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(14), Inches(1))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Voice AI Helpline for Palliative Care"
    p.font.size = Pt(32)
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(14), Inches(0.6))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = '"Companion in Care" — A Digital Public Good for India and Beyond'
    p.font.size = Pt(20)
    p.font.italic = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    # Key features bar
    features = ["15+ Indian Languages", "24/7 Availability", "Clinically Grounded", "Open Source"]
    feature_width = Inches(3.2)
    start_x = Inches(1.2)

    for i, feature in enumerate(features):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      start_x + i * (feature_width + Inches(0.2)),
                                      Inches(7.2), feature_width, Inches(0.7))
        set_shape_color(box, COLORS['white'])
        box.line.fill.background()

        tf = box.text_frame
        tf.paragraphs[0].text = feature
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        set_text_color(tf.paragraphs[0], COLORS['primary'])
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_acknowledgements_slide(prs):
    """Slide 2: Acknowledgements - Sponsors and Partners."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Acknowledgements", color=COLORS['primary'])
    add_subtitle_text(slide, "Funded by Grand Challenges India", top=Inches(1.3), font_size=20, color=COLORS['dark'])

    # Primary funders section
    funders_title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(15), Inches(0.5))
    tf = funders_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Primary Funders"
    p.font.size = Pt(24)
    p.font.bold = True
    set_text_color(p, COLORS['secondary'])
    p.alignment = PP_ALIGN.CENTER

    # Funder boxes
    funders = [
        ("Gates Foundation", "Bill & Melinda Gates Foundation"),
        ("BIRAC-DBT", "Biotechnology Industry Research Assistance Council"),
    ]

    for i, (name, full_name) in enumerate(funders):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(3) + i * Inches(5.5), Inches(2.8),
                                      Inches(4.5), Inches(1.2))
        set_shape_color(box, COLORS['light'])
        set_line_color(box, COLORS['primary'])

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(20)
        p.font.bold = True
        set_text_color(p, COLORS['primary'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = full_name
        p2.font.size = Pt(12)
        set_text_color(p2, COLORS['dark'])
        p2.alignment = PP_ALIGN.CENTER

    # Team section
    team_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(15), Inches(0.5))
    tf = team_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Principal Investigators"
    p.font.size = Pt(24)
    p.font.bold = True
    set_text_color(p, COLORS['secondary'])
    p.alignment = PP_ALIGN.CENTER

    investigators = [
        ("Dr. Anurag Agrawal", "Principal Investigator", "Ashoka University / KCDH"),
        ("Ashish Makani", "Co-Investigator", "inventcures.github.io"),
    ]

    for i, (name, role, affiliation) in enumerate(investigators):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(3) + i * Inches(5.5), Inches(5),
                                      Inches(4.5), Inches(1.3))
        set_shape_color(box, COLORS['white'])
        set_line_color(box, COLORS['secondary'])

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        set_text_color(p, COLORS['dark'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = role
        p2.font.size = Pt(14)
        set_text_color(p2, COLORS['secondary'])
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = affiliation
        p3.font.size = Pt(12)
        set_text_color(p3, COLORS['dark'])
        p3.alignment = PP_ALIGN.CENTER

    # Clinical Partners
    partners_title = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(15), Inches(0.4))
    tf = partners_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Clinical Partners: Pallium India  |  Max Healthcare  |  CanSupport  |  DNipCare"
    p.font.size = Pt(16)
    set_text_color(p, COLORS['dark'])
    p.alignment = PP_ALIGN.CENTER

    # Special thanks
    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(15), Inches(0.8))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "Special thanks to Dr. Parth Sharma (ASAR/Nivarana) for facilitating Pallium India collaboration"
    p.font.size = Pt(14)
    p.font.italic = True
    set_text_color(p, COLORS['dark'])
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Mentorship: Padma Shri Dr. M.R. Rajagopal — Father of Palliative Care in India"
    p2.font.size = Pt(14)
    p2.font.italic = True
    set_text_color(p2, COLORS['primary'])
    p2.alignment = PP_ALIGN.CENTER

def add_crisis_slide(prs):
    """Slide 3: The Palliative Care Crisis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "The Palliative Care Crisis", color=COLORS['red'])

    # Big stat box
    stat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(2), Inches(2), Inches(12), Inches(2.5))
    set_shape_color(stat_box, COLORS['red'])
    stat_box.line.fill.background()

    tf = stat_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Only 1-2%"
    p.font.size = Pt(72)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "of 10+ million Indians needing palliative care receive it"
    p2.font.size = Pt(28)
    set_text_color(p2, COLORS['white'])
    p2.alignment = PP_ALIGN.CENTER

    # Additional stats
    stats = [
        ("57M", "People need palliative\ncare globally"),
        ("78%", "Live in low & middle\nincome countries"),
        ("10M", "Health worker shortage\nprojected by 2030"),
    ]

    for i, (number, description) in enumerate(stats):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(1.5) + i * Inches(4.5), Inches(5.2),
                                      Inches(4), Inches(2))
        set_shape_color(box, COLORS['light'])
        set_line_color(box, COLORS['dark'])

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(48)
        p.font.bold = True
        set_text_color(p, COLORS['red'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = description
        p2.font.size = Pt(16)
        set_text_color(p2, COLORS['dark'])
        p2.alignment = PP_ALIGN.CENTER

    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.8), Inches(15), Inches(0.5))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "98% of patients and their families are left without support"
    p.font.size = Pt(24)
    p.font.bold = True
    set_text_color(p, COLORS['dark'])
    p.alignment = PP_ALIGN.CENTER

def add_opportunity_slide(prs):
    """Slide 4: The Opportunity in India."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "India's Digital Infrastructure: An Opportunity", color=COLORS['primary'])

    # Stats in boxes
    opportunities = [
        ("500M+", "WhatsApp Users", COLORS['primary']),
        ("1.2B", "Mobile Subscribers", COLORS['secondary']),
        ("22", "Scheduled Languages", COLORS['accent']),
        ("Bhashini", "AI Language Platform", COLORS['blue']),
    ]

    for i, (stat, label, color) in enumerate(opportunities):
        row = i // 2
        col = i % 2

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(2) + col * Inches(6),
                                      Inches(2) + row * Inches(2.5),
                                      Inches(5), Inches(2))
        set_shape_color(box, color)
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stat
        p.font.size = Pt(48)
        p.font.bold = True
        set_text_color(p, COLORS['white'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(20)
        set_text_color(p2, COLORS['white'])
        p2.alignment = PP_ALIGN.CENTER

    # Key insight
    insight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(2), Inches(7.2), Inches(12), Inches(1))
    set_shape_color(insight, COLORS['light'])
    set_line_color(insight, COLORS['primary'])

    tf = insight.text_frame
    p = tf.paragraphs[0]
    p.text = "Voice-first AI can reach anyone with a phone — no app, no internet, no literacy required"
    p.font.size = Pt(22)
    p.font.bold = True
    set_text_color(p, COLORS['primary'])
    p.alignment = PP_ALIGN.CENTER

def add_introducing_slide(prs):
    """Slide 5: Introducing Palli Sahayak."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Introducing Palli Sahayak", color=COLORS['primary'])
    add_subtitle_text(slide, "AI Voice Helpline for Palliative Care", top=Inches(1.2), color=COLORS['secondary'])

    # What it does - main features
    features = [
        ("Phone Calls", "Dial in any language, get\nclinically-grounded answers"),
        ("WhatsApp", "Text or voice messages\nin 15+ Indian languages"),
        ("Web Voice", "Real-time voice streaming\nvia browser"),
        ("Always Safe", "Emergency detection\n& professional referral"),
    ]

    for i, (title, desc) in enumerate(features):
        row = i // 2
        col = i % 2

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(1.5) + col * Inches(7),
                                      Inches(2.2) + row * Inches(2.5),
                                      Inches(6), Inches(2))
        set_shape_color(box, COLORS['white'])
        set_line_color(box, COLORS['primary'])

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        set_text_color(p, COLORS['primary'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(16)
        set_text_color(p2, COLORS['dark'])
        p2.alignment = PP_ALIGN.CENTER

    # Key message
    message = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(2), Inches(7.5), Inches(12), Inches(0.8))
    set_shape_color(message, COLORS['primary'])
    message.line.fill.background()

    tf = message.text_frame
    p = tf.paragraphs[0]
    p.text = "Answering questions at 2 AM when families need help most"
    p.font.size = Pt(22)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

def add_architecture_slide(prs):
    """Slide 6: System Architecture Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "System Architecture", color=COLORS['primary'])

    # Architecture diagram using shapes
    # Input channels
    channels = ["Phone\n(Bolna.ai)", "WhatsApp\n(Twilio)", "Web Voice\n(Gemini Live)"]
    for i, channel in enumerate(channels):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(1.5) + i * Inches(4.5), Inches(1.8),
                                      Inches(3.5), Inches(1))
        set_shape_color(box, COLORS['secondary'])
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = channel
        p.font.size = Pt(16)
        p.font.bold = True
        set_text_color(p, COLORS['white'])
        p.alignment = PP_ALIGN.CENTER

    # Arrow down
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(7.5), Inches(2.9), Inches(1), Inches(0.6))
    set_shape_color(arrow1, COLORS['dark'])

    # Voice Router
    router = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(4), Inches(3.6), Inches(8), Inches(0.9))
    set_shape_color(router, COLORS['accent'])
    router.line.fill.background()

    tf = router.text_frame
    p = tf.paragraphs[0]
    p.text = "Voice Router & Smart Query Classifier"
    p.font.size = Pt(20)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    # Arrow down
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(7.5), Inches(4.6), Inches(1), Inches(0.6))
    set_shape_color(arrow2, COLORS['dark'])

    # RAG Pipeline - the core
    rag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(1.5), Inches(5.3), Inches(13), Inches(1.5))
    set_shape_color(rag, COLORS['primary'])
    rag.line.fill.background()

    tf = rag.text_frame
    p = tf.paragraphs[0]
    p.text = "Hybrid RAG Pipeline"
    p.font.size = Pt(24)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "ChromaDB Vector Search  +  Neo4j Knowledge Graph  +  GraphRAG Community Reports"
    p2.font.size = Pt(16)
    set_text_color(p2, COLORS['white'])
    p2.alignment = PP_ALIGN.CENTER

    # Arrow down
    arrow3 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(7.5), Inches(6.9), Inches(1), Inches(0.6))
    set_shape_color(arrow3, COLORS['dark'])

    # Output
    output = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(4), Inches(7.6), Inches(8), Inches(0.9))
    set_shape_color(output, COLORS['blue'])
    output.line.fill.background()

    tf = output.text_frame
    p = tf.paragraphs[0]
    p.text = "Clinically Validated Voice/Text Response"
    p.font.size = Pt(20)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

def add_voice_first_slide(prs):
    """Slide 7: Why Voice-First Design."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Why Voice-First Design?", color=COLORS['primary'])

    # Barriers overcome
    barriers = [
        ("Literacy", "Voice works for everyone\nregardless of reading ability"),
        ("Technology", "No smartphone needed\nJust a basic phone call"),
        ("Language", "Speak in your mother tongue\n15+ Indian languages"),
        ("Geography", "Rural areas with voice\nonly connectivity"),
        ("Timing", "24/7 availability\nHelp at 2 AM"),
        ("Dignity", "Understood in your language\nrespected in your culture"),
    ]

    for i, (title, desc) in enumerate(barriers):
        row = i // 3
        col = i % 3

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8) + col * Inches(5),
                                      Inches(2) + row * Inches(2.8),
                                      Inches(4.5), Inches(2.3))
        set_shape_color(box, COLORS['light'])
        set_line_color(box, COLORS['primary'])

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        set_text_color(p, COLORS['primary'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(16)
        set_text_color(p2, COLORS['dark'])
        p2.alignment = PP_ALIGN.CENTER

    # Bottom message
    message = slide.shapes.add_textbox(Inches(0.5), Inches(7.8), Inches(15), Inches(0.5))
    tf = message.text_frame
    p = tf.paragraphs[0]
    p.text = "Voice AI democratizes access to quality healthcare information"
    p.font.size = Pt(22)
    p.font.bold = True
    set_text_color(p, COLORS['secondary'])
    p.alignment = PP_ALIGN.CENTER

def add_languages_slide(prs):
    """Slide 8: Language Support."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "15+ Indian Languages", color=COLORS['primary'])
    add_subtitle_text(slide, "Speaking to users in their mother tongue", top=Inches(1.2), color=COLORS['secondary'])

    # Language grid
    languages = [
        ("Hindi", "हिन्दी", "P0"),
        ("English", "English", "P0"),
        ("Bengali", "বাংলা", "P0"),
        ("Tamil", "தமிழ்", "P0"),
        ("Telugu", "తెలుగు", "P1"),
        ("Marathi", "मराठी", "P1"),
        ("Gujarati", "ગુજરાતી", "P1"),
        ("Kannada", "ಕನ್ನಡ", "P1"),
        ("Malayalam", "മലയാളം", "P2"),
        ("Punjabi", "ਪੰਜਾਬੀ", "P2"),
        ("Odia", "ଓଡ଼ିଆ", "P2"),
        ("Assamese", "অসমীয়া", "P3"),
    ]

    for i, (name, script, priority) in enumerate(languages):
        row = i // 4
        col = i % 4

        color = COLORS['primary'] if priority == "P0" else (
            COLORS['secondary'] if priority == "P1" else COLORS['accent']
        )

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8) + col * Inches(3.8),
                                      Inches(2.2) + row * Inches(1.6),
                                      Inches(3.4), Inches(1.3))
        set_shape_color(box, color)
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = script
        p.font.size = Pt(22)
        p.font.bold = True
        set_text_color(p, COLORS['white'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = name
        p2.font.size = Pt(14)
        set_text_color(p2, COLORS['white'])
        p2.alignment = PP_ALIGN.CENTER

    # Bhashini integration note
    note = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(2), Inches(7.2), Inches(12), Inches(1))
    set_shape_color(note, COLORS['light'])
    set_line_color(note, COLORS['blue'])

    tf = note.text_frame
    p = tf.paragraphs[0]
    p.text = "Powered by Bhashini — India's Government AI Language Platform"
    p.font.size = Pt(20)
    p.font.bold = True
    set_text_color(p, COLORS['blue'])
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Free API access for Digital Public Goods | 22 scheduled languages supported"
    p2.font.size = Pt(14)
    set_text_color(p2, COLORS['dark'])
    p2.alignment = PP_ALIGN.CENTER

def add_rag_foundation_slide(prs):
    """Slide 9: Clinical RAG Foundation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Clinical RAG Foundation", color=COLORS['primary'])
    add_subtitle_text(slide, "Retrieval-Augmented Generation with Medical Guidelines", top=Inches(1.2), color=COLORS['secondary'])

    # RAG explanation diagram
    steps = [
        ("1. User Question", "Voice or text query\nin any language", COLORS['secondary']),
        ("2. Query Processing", "Translation + entity\nextraction", COLORS['accent']),
        ("3. Knowledge Retrieval", "Search clinical\nguidelines", COLORS['primary']),
        ("4. Response Generation", "Contextual answer\nwith citations", COLORS['blue']),
    ]

    for i, (title, desc, color) in enumerate(steps):
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8) + i * Inches(3.8), Inches(2.3),
                                      Inches(3.4), Inches(2))
        set_shape_color(box, color)
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        set_text_color(p, COLORS['white'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        set_text_color(p2, COLORS['white'])
        p2.alignment = PP_ALIGN.CENTER

        # Arrow (except last)
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                            Inches(4) + i * Inches(3.8), Inches(3),
                                            Inches(0.5), Inches(0.4))
            set_shape_color(arrow, COLORS['dark'])

    # Source documents
    sources_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(15), Inches(0.5))
    tf = sources_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Grounded in Clinical Guidelines"
    p.font.size = Pt(22)
    p.font.bold = True
    set_text_color(p, COLORS['primary'])
    p.alignment = PP_ALIGN.CENTER

    sources = [
        "WHO Palliative Care",
        "IAHPC Essential Practices",
        "IAPC Handbook (India)",
        "Pallium India Protocols",
    ]

    for i, source in enumerate(sources):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8) + i * Inches(3.8), Inches(5.5),
                                      Inches(3.4), Inches(0.8))
        set_shape_color(box, COLORS['light'])
        set_line_color(box, COLORS['primary'])

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = source
        p.font.size = Pt(14)
        p.font.bold = True
        set_text_color(p, COLORS['dark'])
        p.alignment = PP_ALIGN.CENTER

    # Key benefit
    benefit = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(2), Inches(6.8), Inches(12), Inches(1.5))
    set_shape_color(benefit, COLORS['primary'])
    benefit.line.fill.background()

    tf = benefit.text_frame
    p = tf.paragraphs[0]
    p.text = "Every response is grounded in verified medical sources"
    p.font.size = Pt(22)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Citations provided | No hallucinations | Expert-validated content"
    p2.font.size = Pt(16)
    set_text_color(p2, COLORS['white'])
    p2.alignment = PP_ALIGN.CENTER

def add_knowledge_architecture_slide(prs):
    """Slide 10: Triple-Layer Knowledge Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Triple-Layer Knowledge Architecture", color=COLORS['primary'])

    # Three layers
    layers = [
        ("ChromaDB", "Vector Search", "Semantic similarity\n384-dimensional embeddings", COLORS['secondary']),
        ("Neo4j", "Knowledge Graph", "Structured relationships\n18 entity types", COLORS['primary']),
        ("GraphRAG", "Community Reports", "Multi-hop reasoning\nLeiden clustering", COLORS['blue']),
    ]

    for i, (name, subtitle, desc, color) in enumerate(layers):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8) + i * Inches(5), Inches(2),
                                      Inches(4.5), Inches(3.5))
        set_shape_color(box, color)
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(32)
        p.font.bold = True
        set_text_color(p, COLORS['white'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        set_text_color(p2, COLORS['white'])
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = ""  # spacer
        p3.font.size = Pt(8)

        p4 = tf.add_paragraph()
        p4.text = desc
        p4.font.size = Pt(14)
        set_text_color(p4, COLORS['white'])
        p4.alignment = PP_ALIGN.CENTER

    # Fusion strategy
    fusion = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(2), Inches(6), Inches(12), Inches(2.3))
    set_shape_color(fusion, COLORS['light'])
    set_line_color(fusion, COLORS['accent'])

    tf = fusion.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Intelligent Context Fusion"
    p.font.size = Pt(24)
    p.font.bold = True
    set_text_color(p, COLORS['accent'])
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "FUSE: When multiple contexts are equally relevant (distance < 0.15)"
    p2.font.size = Pt(16)
    set_text_color(p2, COLORS['dark'])
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "FOCUS: When one context is clearly best"
    p3.font.size = Pt(16)
    set_text_color(p3, COLORS['dark'])
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = "Prevents hallucination while maximizing relevant context"
    p4.font.size = Pt(14)
    p4.font.italic = True
    set_text_color(p4, COLORS['secondary'])
    p4.alignment = PP_ALIGN.CENTER

def add_clinical_validation_slide(prs):
    """Slide 11: Clinical Validation Pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Clinical Validation Pipeline", color=COLORS['primary'])
    add_subtitle_text(slide, "Ensuring Safe, Accurate Medical Information", top=Inches(1.2), color=COLORS['red'])

    # Validation steps
    steps = [
        ("Dosage Check", "Validate medication\ndoses against ranges"),
        ("Safety Keywords", "Detect emergencies\nand urgent situations"),
        ("Hallucination Check", "Verify response against\nsource documents"),
        ("Citation Verify", "Ensure proper source\nattribution"),
    ]

    for i, (title, desc) in enumerate(steps):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8) + i * Inches(3.8), Inches(2.2),
                                      Inches(3.4), Inches(2))
        set_shape_color(box, COLORS['light'])
        set_line_color(box, COLORS['primary'])

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        set_text_color(p, COLORS['primary'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        set_text_color(p2, COLORS['dark'])
        p2.alignment = PP_ALIGN.CENTER

    # Expert sampling
    sampling_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(15), Inches(0.5))
    tf = sampling_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Expert Review Sampling"
    p.font.size = Pt(22)
    p.font.bold = True
    set_text_color(p, COLORS['secondary'])
    p.alignment = PP_ALIGN.CENTER

    sampling = [
        ("5%", "Normal queries", COLORS['primary']),
        ("50%", "High-risk topics", COLORS['accent']),
        ("100%", "Critical/Emergency", COLORS['red']),
    ]

    for i, (pct, label, color) in enumerate(sampling):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(2.5) + i * Inches(4), Inches(5.2),
                                      Inches(3), Inches(1.3))
        set_shape_color(box, color)
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = pct
        p.font.size = Pt(36)
        p.font.bold = True
        set_text_color(p, COLORS['white'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(14)
        set_text_color(p2, COLORS['white'])
        p2.alignment = PP_ALIGN.CENTER

    # Target metrics
    metrics = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(2), Inches(7), Inches(12), Inches(1.3))
    set_shape_color(metrics, COLORS['primary'])
    metrics.line.fill.background()

    tf = metrics.text_frame
    p = tf.paragraphs[0]
    p.text = "Target: 95%+ Accuracy  |  <2% Hallucination  |  90%+ Expert Agreement"
    p.font.size = Pt(20)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

def add_personalization_slide(prs):
    """Slide 12: User Personalization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Adaptive User Experience", color=COLORS['primary'])

    # User roles
    roles_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(15), Inches(0.5))
    tf = roles_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Automatic Role Detection"
    p.font.size = Pt(22)
    p.font.bold = True
    set_text_color(p, COLORS['secondary'])
    p.alignment = PP_ALIGN.CENTER

    roles = [
        ("Patient", "Personal symptom\nmanagement", '"मुझे दर्द हो रहा है"'),
        ("Caregiver", "Care guidance for\nfamily member", '"my mother has..."'),
        ("Healthcare Worker", "Clinical decision\nsupport", '"my patient needs..."'),
    ]

    for i, (role, desc, example) in enumerate(roles):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(1) + i * Inches(4.8), Inches(2.2),
                                      Inches(4.3), Inches(2.2))
        set_shape_color(box, COLORS['light'])
        set_line_color(box, COLORS['primary'])

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = role
        p.font.size = Pt(22)
        p.font.bold = True
        set_text_color(p, COLORS['primary'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        set_text_color(p2, COLORS['dark'])
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = example
        p3.font.size = Pt(12)
        p3.font.italic = True
        set_text_color(p3, COLORS['secondary'])
        p3.alignment = PP_ALIGN.CENTER

    # Context memory
    memory_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(15), Inches(0.5))
    tf = memory_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Context Memory Across Sessions"
    p.font.size = Pt(22)
    p.font.bold = True
    set_text_color(p, COLORS['secondary'])
    p.alignment = PP_ALIGN.CENTER

    memory_items = [
        ("Patient Condition", "Cancer, COPD, Heart Failure"),
        ("Current Symptoms", "Pain, Nausea, Fatigue"),
        ("Medications", "Morphine 10mg, Ondansetron"),
        ("Allergies", "Drug sensitivities recorded"),
    ]

    for i, (title, example) in enumerate(memory_items):
        row = i // 2
        col = i % 2

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(2) + col * Inches(6),
                                      Inches(5.3) + row * Inches(1.2),
                                      Inches(5.5), Inches(1))
        set_shape_color(box, COLORS['white'])
        set_line_color(box, COLORS['secondary'])

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{title}: {example}"
        p.font.size = Pt(14)
        set_text_color(p, COLORS['dark'])
        p.alignment = PP_ALIGN.CENTER

    # Key benefit
    benefit = slide.shapes.add_textbox(Inches(0.5), Inches(8), Inches(15), Inches(0.5))
    tf = benefit.text_frame
    p = tf.paragraphs[0]
    p.text = '"Last time you asked about pain management..." — Continuity of care'
    p.font.size = Pt(18)
    p.font.italic = True
    set_text_color(p, COLORS['primary'])
    p.alignment = PP_ALIGN.CENTER

def add_analytics_slide(prs):
    """Slide 13: Real-time Analytics Dashboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Real-time Analytics Dashboard", color=COLORS['primary'])

    # Metrics grid
    metrics = [
        ("Response Time", "P50: 400ms\nP95: 800ms\nP99: 1.2s", COLORS['primary']),
        ("RAG Success", "89.3%\nretrieval rate", COLORS['secondary']),
        ("Validation Rate", "95.2%\npass rate", COLORS['primary']),
        ("User Satisfaction", "4.6/5.0\naverage", COLORS['accent']),
    ]

    for i, (title, value, color) in enumerate(metrics):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8) + i * Inches(3.8), Inches(1.8),
                                      Inches(3.4), Inches(1.8))
        set_shape_color(box, color)
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        set_text_color(p, COLORS['white'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(16)
        set_text_color(p2, COLORS['white'])
        p2.alignment = PP_ALIGN.CENTER

    # Language distribution (mock bar chart with shapes)
    lang_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(7), Inches(0.5))
    tf = lang_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Language Distribution"
    p.font.size = Pt(20)
    p.font.bold = True
    set_text_color(p, COLORS['dark'])

    languages = [("Hindi", 45), ("Bengali", 18), ("Tamil", 12), ("English", 10), ("Others", 15)]
    max_width = Inches(5)

    for i, (lang, pct) in enumerate(languages):
        # Bar
        bar_width = max_width * (pct / 50)  # Scale to max 50%
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(1.5), Inches(4.5) + i * Inches(0.5),
                                      bar_width, Inches(0.35))
        set_shape_color(bar, COLORS['primary'])
        bar.line.fill.background()

        # Label
        label = slide.shapes.add_textbox(Inches(0.5), Inches(4.5) + i * Inches(0.5), Inches(1), Inches(0.35))
        tf = label.text_frame
        p = tf.paragraphs[0]
        p.text = lang
        p.font.size = Pt(12)
        set_text_color(p, COLORS['dark'])

        # Percentage
        pct_label = slide.shapes.add_textbox(Inches(1.5) + bar_width + Inches(0.1),
                                              Inches(4.5) + i * Inches(0.5),
                                              Inches(0.5), Inches(0.35))
        tf = pct_label.text_frame
        p = tf.paragraphs[0]
        p.text = f"{pct}%"
        p.font.size = Pt(12)
        p.font.bold = True
        set_text_color(p, COLORS['dark'])

    # Topic distribution
    topic_title = slide.shapes.add_textbox(Inches(8), Inches(3.9), Inches(7), Inches(0.5))
    tf = topic_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Topic Distribution"
    p.font.size = Pt(20)
    p.font.bold = True
    set_text_color(p, COLORS['dark'])

    topics = [("Pain Mgmt", 32), ("Symptoms", 24), ("Medications", 18), ("Caregiver", 14), ("End-of-Life", 12)]

    for i, (topic, pct) in enumerate(topics):
        bar_width = max_width * (pct / 35)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(9.5), Inches(4.5) + i * Inches(0.5),
                                      bar_width, Inches(0.35))
        set_shape_color(bar, COLORS['secondary'])
        bar.line.fill.background()

        label = slide.shapes.add_textbox(Inches(8), Inches(4.5) + i * Inches(0.5), Inches(1.5), Inches(0.35))
        tf = label.text_frame
        p = tf.paragraphs[0]
        p.text = topic
        p.font.size = Pt(12)
        set_text_color(p, COLORS['dark'])

        pct_label = slide.shapes.add_textbox(Inches(9.5) + bar_width + Inches(0.1),
                                              Inches(4.5) + i * Inches(0.5),
                                              Inches(0.5), Inches(0.35))
        tf = pct_label.text_frame
        p = tf.paragraphs[0]
        p.text = f"{pct}%"
        p.font.size = Pt(12)
        p.font.bold = True
        set_text_color(p, COLORS['dark'])

    # Health status indicator
    health = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(2), Inches(7.5), Inches(12), Inches(0.8))
    set_shape_color(health, COLORS['primary'])
    health.line.fill.background()

    tf = health.text_frame
    p = tf.paragraphs[0]
    p.text = "System Health: HEALTHY  |  Alerts: 0  |  Uptime: 99.9%"
    p.font.size = Pt(18)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

def add_dpg_slide(prs):
    """Slide 14: Digital Public Good Positioning."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "A Digital Public Good (DPG)", color=COLORS['primary'])
    add_subtitle_text(slide, "Open Source | Replicable | Sustainable", top=Inches(1.2), color=COLORS['secondary'])

    # DPG criteria
    criteria = [
        ("Open Source", "MIT License\nPublicly on GitHub"),
        ("SDG Aligned", "SDG 3: Good Health\nSDG 10: Reduced Inequalities"),
        ("Privacy First", "No personal health data\nAnonymized interactions"),
        ("Do No Harm", "Medical guardrails\nProfessional referral"),
        ("Interoperable", "REST APIs\nModular architecture"),
        ("Data Standards", "SNOMED, ICD-10\nMedical ontologies"),
    ]

    for i, (title, desc) in enumerate(criteria):
        row = i // 3
        col = i % 3

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8) + col * Inches(5),
                                      Inches(2.2) + row * Inches(2.3),
                                      Inches(4.5), Inches(2))
        set_shape_color(box, COLORS['light'])
        set_line_color(box, COLORS['primary'])

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"✓ {title}"
        p.font.size = Pt(20)
        p.font.bold = True
        set_text_color(p, COLORS['primary'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        set_text_color(p2, COLORS['dark'])
        p2.alignment = PP_ALIGN.CENTER

    # DPGA alignment
    dpga = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(2), Inches(7.2), Inches(12), Inches(1))
    set_shape_color(dpga, COLORS['blue'])
    dpga.line.fill.background()

    tf = dpga.text_frame
    p = tf.paragraphs[0]
    p.text = "Aligned with Digital Public Goods Alliance (DPGA) Standards"
    p.font.size = Pt(20)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

def add_global_replicability_slide(prs):
    """Slide 15: Global Replicability."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Beyond India: Global Replicability", color=COLORS['primary'])
    add_subtitle_text(slide, "A model for palliative care AI worldwide", top=Inches(1.2), color=COLORS['secondary'])

    # Configuration adaptability
    config_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.8), Inches(2), Inches(6.5), Inches(4))
    set_shape_color(config_box, COLORS['dark'])
    config_box.line.fill.background()

    tf = config_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "# config.yaml - Kenya Example"
    p.font.size = Pt(14)
    set_text_color(p, COLORS['accent'])
    p.alignment = PP_ALIGN.LEFT

    config_text = """
country: "kenya"
languages:
  - code: "sw"
    name: "Swahili"
  - code: "en"
    name: "English"

corpus:
  source: "kenya_palliative_guidelines"
"""

    p2 = tf.add_paragraph()
    p2.text = config_text
    p2.font.size = Pt(12)
    set_text_color(p2, COLORS['white'])
    p2.alignment = PP_ALIGN.LEFT

    # Potential regions
    regions_title = slide.shapes.add_textbox(Inches(8), Inches(2), Inches(7), Inches(0.5))
    tf = regions_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Potential Adaptations"
    p.font.size = Pt(22)
    p.font.bold = True
    set_text_color(p, COLORS['secondary'])

    regions = [
        ("Sub-Saharan Africa", "Kenya, Nigeria, South Africa"),
        ("Southeast Asia", "Indonesia, Philippines, Vietnam"),
        ("Latin America", "Brazil, Mexico, Colombia"),
        ("Middle East", "Egypt, Jordan, Lebanon"),
        ("Developed Nations", "Rural US, UK NHS, Australia"),
    ]

    for i, (region, countries) in enumerate(regions):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(8), Inches(2.7) + i * Inches(1.1),
                                      Inches(7), Inches(0.9))
        set_shape_color(box, COLORS['light'])
        set_line_color(box, COLORS['primary'])

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{region}"
        p.font.size = Pt(16)
        p.font.bold = True
        set_text_color(p, COLORS['primary'])
        p.alignment = PP_ALIGN.LEFT

        p2 = tf.add_paragraph()
        p2.text = countries
        p2.font.size = Pt(12)
        set_text_color(p2, COLORS['dark'])
        p2.alignment = PP_ALIGN.LEFT

    # Partnership opportunities
    partnerships = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(0.8), Inches(6.5), Inches(14.4), Inches(1.8))
    set_shape_color(partnerships, COLORS['primary'])
    partnerships.line.fill.background()

    tf = partnerships.text_frame
    p = tf.paragraphs[0]
    p.text = "Partnership Opportunities"
    p.font.size = Pt(22)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "WHO | UNICEF | Open Source Pharma | Local Health Ministries | IAHPC"
    p2.font.size = Pt(18)
    set_text_color(p2, COLORS['white'])
    p2.alignment = PP_ALIGN.CENTER

def add_impact_metrics_slide(prs):
    """Slide 16: Impact Metrics & Targets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Impact Metrics & Targets", color=COLORS['primary'])

    # V2 targets table
    targets = [
        ("Monthly Active Users", "1,000", "100,000", "100x"),
        ("Languages Supported", "6", "15+", "2.5x"),
        ("Response Accuracy", "85%", "95%+", "+10%"),
        ("Voice Latency", "2.5s", "<0.5s", "5x faster"),
        ("Healthcare Workers", "100", "10,000", "100x"),
        ("Cost per Query", "$0.02", "$0.005", "4x cheaper"),
    ]

    # Header
    headers = ["Metric", "V1 Baseline", "V2 Target", "Improvement"]
    header_widths = [Inches(4), Inches(2.5), Inches(2.5), Inches(2.5)]

    for i, (header, width) in enumerate(zip(headers, header_widths)):
        x_pos = Inches(1.5) + sum(hw for hw in header_widths[:i])
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x_pos, Inches(1.8), width, Inches(0.6))
        set_shape_color(box, COLORS['primary'])
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(16)
        p.font.bold = True
        set_text_color(p, COLORS['white'])
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, (metric, v1, v2, improvement) in enumerate(targets):
        values = [metric, v1, v2, improvement]
        row_color = COLORS['light'] if row_idx % 2 == 0 else COLORS['white']

        for col_idx, (value, width) in enumerate(zip(values, header_widths)):
            x_pos = Inches(1.5) + sum(hw for hw in header_widths[:col_idx])
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                          x_pos, Inches(2.4) + row_idx * Inches(0.7),
                                          width, Inches(0.7))
            set_shape_color(box, row_color)
            set_line_color(box, COLORS['dark'])

            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = value
            p.font.size = Pt(14)
            if col_idx == 3:  # Improvement column
                p.font.bold = True
                set_text_color(p, COLORS['primary'])
            else:
                set_text_color(p, COLORS['dark'])
            p.alignment = PP_ALIGN.CENTER

    # Long-term vision
    vision = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(1.5), Inches(6.8), Inches(13), Inches(1.5))
    set_shape_color(vision, COLORS['secondary'])
    vision.line.fill.background()

    tf = vision.text_frame
    p = tf.paragraphs[0]
    p.text = "2-Year Vision"
    p.font.size = Pt(22)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "1M Users  |  500K Care Decisions/Year  |  10 Countries  |  5 State Adoptions"
    p2.font.size = Pt(18)
    set_text_color(p2, COLORS['white'])
    p2.alignment = PP_ALIGN.CENTER

def add_partners_slide(prs):
    """Slide 17: Partners & Collaborations."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Partners & Collaborations", color=COLORS['primary'])

    # Partner categories
    categories = [
        ("Clinical Partners", ["Pallium India", "Max Healthcare", "CanSupport", "DNipCare"], COLORS['primary']),
        ("Academic", ["Ashoka University", "KCDH", "ASAR/Nivarana"], COLORS['secondary']),
        ("Government", ["BIRAC-DBT", "Grand Challenges India", "Bhashini"], COLORS['blue']),
        ("Funders", ["Gates Foundation", "BIRAC"], COLORS['accent']),
    ]

    for i, (category, partners, color) in enumerate(categories):
        row = i // 2
        col = i % 2

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(1) + col * Inches(7.5),
                                      Inches(1.8) + row * Inches(2.8),
                                      Inches(6.5), Inches(2.5))
        set_shape_color(box, COLORS['light'])
        set_line_color(box, color)

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category
        p.font.size = Pt(22)
        p.font.bold = True
        set_text_color(p, color)
        p.alignment = PP_ALIGN.CENTER

        for partner in partners:
            p2 = tf.add_paragraph()
            p2.text = f"• {partner}"
            p2.font.size = Pt(16)
            set_text_color(p2, COLORS['dark'])
            p2.alignment = PP_ALIGN.CENTER

    # Special acknowledgement
    ack = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(1), Inches(7.2), Inches(14), Inches(1))
    set_shape_color(ack, COLORS['primary'])
    ack.line.fill.background()

    tf = ack.text_frame
    p = tf.paragraphs[0]
    p.text = "Mentorship: Padma Shri Dr. M.R. Rajagopal — Father of Palliative Care in India"
    p.font.size = Pt(18)
    p.font.italic = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

def add_tech_stack_slide(prs):
    """Slide 18: Technology Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Technology Stack", color=COLORS['primary'])

    # Tech categories
    tech = [
        ("Framework", "FastAPI • Uvicorn • Gradio Admin", COLORS['blue']),
        ("Voice AI", "Bolna.ai • Gemini Live • Groq Whisper • Edge TTS", COLORS['secondary']),
        ("RAG Pipeline", "ChromaDB • Neo4j • Microsoft GraphRAG", COLORS['primary']),
        ("LLMs", "Groq (llama-3.1-8b) • Qwen3-32B • GPT-4o-mini", COLORS['accent']),
        ("Languages", "Bhashini • Deepgram • ElevenLabs", COLORS['secondary']),
        ("Infrastructure", "Python 3.11 • Docker • ngrok", COLORS['dark']),
    ]

    for i, (category, items, color) in enumerate(tech):
        row = i // 2
        col = i % 2

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8) + col * Inches(7.5),
                                      Inches(1.8) + row * Inches(2),
                                      Inches(7), Inches(1.7))
        set_shape_color(box, color)
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category
        p.font.size = Pt(20)
        p.font.bold = True
        set_text_color(p, COLORS['white'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = items
        p2.font.size = Pt(14)
        set_text_color(p2, COLORS['white'])
        p2.alignment = PP_ALIGN.CENTER

    # Open source note
    opensource = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(2), Inches(7.5), Inches(12), Inches(0.8))
    set_shape_color(opensource, COLORS['light'])
    set_line_color(opensource, COLORS['primary'])

    tf = opensource.text_frame
    p = tf.paragraphs[0]
    p.text = "100% Open Source • github.com/inventcures/rag_gci • MIT License"
    p.font.size = Pt(18)
    p.font.bold = True
    set_text_color(p, COLORS['primary'])
    p.alignment = PP_ALIGN.CENTER

def add_indian_context_slide(prs):
    """Slide 19: Indian Socio-Cultural Context."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Situated in Indian Socio-Cultural Context", color=COLORS['primary'])

    # Cultural considerations
    considerations = [
        ("Family-Centered Care", "Healthcare decisions are family matters.\nPalli Sahayak supports family caregivers."),
        ("Spiritual Dimensions", "End-of-life care integrates spiritual practices.\nRespectful of diverse beliefs."),
        ("Home-Based Care", "Most palliative care happens at home.\nDesigned for caregivers, not just professionals."),
        ("Language Diversity", "India's linguistic diversity is respected.\n15+ languages with cultural responses."),
    ]

    for i, (title, desc) in enumerate(considerations):
        row = i // 2
        col = i % 2

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8) + col * Inches(7.5),
                                      Inches(1.8) + row * Inches(2.5),
                                      Inches(7), Inches(2.2))
        set_shape_color(box, COLORS['light'])
        set_line_color(box, COLORS['primary'])

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        set_text_color(p, COLORS['primary'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        set_text_color(p2, COLORS['dark'])
        p2.alignment = PP_ALIGN.CENTER

    # India-specific dataset
    dataset = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(1), Inches(6.8), Inches(14), Inches(1.5))
    set_shape_color(dataset, COLORS['secondary'])
    dataset.line.fill.background()

    tf = dataset.text_frame
    p = tf.paragraphs[0]
    p.text = "India-Specific Case Scenarios Dataset"
    p.font.size = Pt(22)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Curated with Pallium India & Max Healthcare — Reflecting local clinical contexts"
    p2.font.size = Pt(16)
    set_text_color(p2, COLORS['white'])
    p2.alignment = PP_ALIGN.CENTER

def add_roadmap_slide(prs):
    """Slide 20: Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title_text(slide, "Development Roadmap", color=COLORS['primary'])

    # Timeline
    phases = [
        ("Q1-Q2 2026", "Scale & Validate",
         "• GraphRAG integration\n• Clinical validation\n• 15 language rollout\n• Sub-500ms latency", COLORS['primary']),
        ("Q3-Q4 2026", "Expand & Partner",
         "• 100K user target\n• 5+ hospice integrations\n• Research publications\n• State pilot programs", COLORS['secondary']),
        ("2027+", "Global Impact",
         "• 1M users target\n• 10 country adaptations\n• Policy integration\n• Research platform", COLORS['accent']),
    ]

    for i, (timeline, title, items, color) in enumerate(phases):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8) + i * Inches(5), Inches(2),
                                      Inches(4.5), Inches(4.5))
        set_shape_color(box, color)
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = timeline
        p.font.size = Pt(18)
        set_text_color(p, COLORS['white'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(26)
        p2.font.bold = True
        set_text_color(p2, COLORS['white'])
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = items
        p3.font.size = Pt(14)
        set_text_color(p3, COLORS['white'])
        p3.alignment = PP_ALIGN.LEFT

    # Current status
    status = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(2), Inches(7), Inches(12), Inches(1.3))
    set_shape_color(status, COLORS['light'])
    set_line_color(status, COLORS['primary'])

    tf = status.text_frame
    p = tf.paragraphs[0]
    p.text = "Current Status: V2 Development — Clinical Validation, Analytics, Personalization"
    p.font.size = Pt(18)
    p.font.bold = True
    set_text_color(p, COLORS['primary'])
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "10 V2 components tested and verified (100% pass rate)"
    p2.font.size = Pt(14)
    set_text_color(p2, COLORS['dark'])
    p2.alignment = PP_ALIGN.CENTER

def add_call_to_action_slide(prs):
    """Slide 21: Call to Action."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
    set_shape_color(background, COLORS['primary'])
    background.line.fill.background()

    add_title_text(slide, "How You Can Help", top=Inches(0.8), color=COLORS['white'])

    # Action items
    actions = [
        ("Healthcare Professionals", "• Validate clinical content\n• Provide feedback\n• Refer patients and caregivers"),
        ("Researchers", "• Collaborate on effectiveness studies\n• Contribute to datasets\n• Co-author publications"),
        ("Technologists", "• Contribute to open source code\n• Adapt for your region\n• Build integrations"),
        ("Funders & Partners", "• Support infrastructure scaling\n• Fund regional adaptations\n• Policy advocacy"),
    ]

    for i, (audience, items) in enumerate(actions):
        row = i // 2
        col = i % 2

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8) + col * Inches(7.5),
                                      Inches(2) + row * Inches(2.5),
                                      Inches(7), Inches(2.2))
        set_shape_color(box, COLORS['white'])
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = audience
        p.font.size = Pt(22)
        p.font.bold = True
        set_text_color(p, COLORS['primary'])
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = items
        p2.font.size = Pt(14)
        set_text_color(p2, COLORS['dark'])
        p2.alignment = PP_ALIGN.LEFT

    # Links
    links = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(2), Inches(7.2), Inches(12), Inches(1))
    set_shape_color(links, COLORS['white'])
    links.line.fill.background()

    tf = links.text_frame
    p = tf.paragraphs[0]
    p.text = "GitHub: github.com/inventcures/rag_gci  |  Docs: deepwiki.com/inventcures/rag_gci"
    p.font.size = Pt(18)
    p.font.bold = True
    set_text_color(p, COLORS['primary'])
    p.alignment = PP_ALIGN.CENTER

def add_thank_you_slide(prs):
    """Slide 22: Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
    set_shape_color(background, COLORS['primary'])
    background.line.fill.background()

    # Main message
    title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(72)
    p.font.bold = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(0.8))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = "Palli Sahayak — Compassionate AI for Palliative Care"
    p.font.size = Pt(28)
    p.font.italic = True
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER

    # Contact info box
    contact = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(3), Inches(4.8), Inches(10), Inches(2.5))
    set_shape_color(contact, COLORS['white'])
    contact.line.fill.background()

    tf = contact.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Contact & Resources"
    p.font.size = Pt(24)
    p.font.bold = True
    set_text_color(p, COLORS['primary'])
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "GitHub: github.com/inventcures/rag_gci"
    p2.font.size = Pt(16)
    set_text_color(p2, COLORS['dark'])
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Documentation: deepwiki.com/inventcures/rag_gci"
    p3.font.size = Pt(16)
    set_text_color(p3, COLORS['dark'])
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = "Grant: gcgh.grandchallenges.org"
    p4.font.size = Pt(16)
    set_text_color(p4, COLORS['dark'])
    p4.alignment = PP_ALIGN.CENTER

    # Footer
    footer = slide.shapes.add_textbox(Inches(1), Inches(7.8), Inches(14), Inches(0.5))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Open Source • MIT License • A Digital Public Good for India and the World"
    p.font.size = Pt(16)
    set_text_color(p, COLORS['white'])
    p.alignment = PP_ALIGN.CENTER


def create_pdf_presentation():
    """Create PDF version of the presentation using reportlab."""
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT

    output_dir = os.path.dirname(os.path.abspath(__file__))
    pdf_path = os.path.join(output_dir, "Palli_Sahayak_Presentation.pdf")

    # Colors
    PRIMARY = HexColor('#2E7D32')
    SECONDARY = HexColor('#00696B')
    ACCENT = HexColor('#FF8F00')
    DARK = HexColor('#212121')
    RED = HexColor('#C62828')
    BLUE = HexColor('#1565C0')
    WHITE = HexColor('#FFFFFF')
    LIGHT = HexColor('#F5F5F5')

    # Create document
    doc = SimpleDocTemplate(
        pdf_path,
        pagesize=landscape(letter),
        rightMargin=0.5*inch,
        leftMargin=0.5*inch,
        topMargin=0.5*inch,
        bottomMargin=0.5*inch
    )

    # Styles
    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=36,
        textColor=PRIMARY,
        alignment=TA_CENTER,
        spaceAfter=20,
        fontName='Helvetica-Bold'
    )

    subtitle_style = ParagraphStyle(
        'CustomSubtitle',
        parent=styles['Heading2'],
        fontSize=18,
        textColor=SECONDARY,
        alignment=TA_CENTER,
        spaceAfter=15,
        fontName='Helvetica'
    )

    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['Normal'],
        fontSize=12,
        textColor=DARK,
        alignment=TA_LEFT,
        spaceAfter=10,
        fontName='Helvetica'
    )

    center_body = ParagraphStyle(
        'CenterBody',
        parent=body_style,
        alignment=TA_CENTER
    )

    big_stat = ParagraphStyle(
        'BigStat',
        parent=styles['Heading1'],
        fontSize=48,
        textColor=RED,
        alignment=TA_CENTER,
        spaceAfter=10,
        fontName='Helvetica-Bold'
    )

    # Build content
    story = []

    # Slide 1: Title
    story.append(Spacer(1, 1*inch))
    story.append(Paragraph("Palli Sahayak", ParagraphStyle('BigTitle', parent=title_style, fontSize=48, textColor=PRIMARY)))
    story.append(Paragraph("पल्ली सहायक", ParagraphStyle('Hindi', parent=subtitle_style, fontSize=24)))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("Voice AI Helpline for Palliative Care", subtitle_style))
    story.append(Spacer(1, 0.2*inch))
    story.append(Paragraph('"Companion in Care" — A Digital Public Good for India and Beyond', ParagraphStyle('Tagline', parent=body_style, alignment=TA_CENTER, fontName='Helvetica-Oblique')))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph("15+ Indian Languages  |  24/7 Availability  |  Clinically Grounded  |  Open Source", center_body))
    story.append(Spacer(1, 2*inch))

    # Page break
    from reportlab.platypus import PageBreak
    story.append(PageBreak())

    # Slide 2: Acknowledgements
    story.append(Paragraph("Acknowledgements", title_style))
    story.append(Paragraph("Funded by Grand Challenges India", subtitle_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Primary Funders</b>", center_body))
    story.append(Paragraph("• Gates Foundation (Bill & Melinda Gates Foundation)", center_body))
    story.append(Paragraph("• BIRAC-DBT (Biotechnology Industry Research Assistance Council)", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Principal Investigators</b>", center_body))
    story.append(Paragraph("• Dr. Anurag Agrawal — Principal Investigator (Ashoka University / KCDH)", center_body))
    story.append(Paragraph("• Ashish Makani — Co-Investigator (inventcures.github.io)", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Clinical Partners</b>: Pallium India | Max Healthcare | CanSupport | DNipCare", center_body))
    story.append(Spacer(1, 0.2*inch))
    story.append(Paragraph("<i>Mentorship: Padma Shri Dr. M.R. Rajagopal — Father of Palliative Care in India</i>", center_body))
    story.append(PageBreak())

    # Slide 3: Crisis
    story.append(Paragraph("The Palliative Care Crisis", ParagraphStyle('Crisis', parent=title_style, textColor=RED)))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("Only 1-2%", big_stat))
    story.append(Paragraph("of 10+ million Indians needing palliative care receive it", center_body))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph("<b>57M</b> people globally need palliative care", center_body))
    story.append(Paragraph("<b>78%</b> live in low & middle income countries", center_body))
    story.append(Paragraph("<b>10M</b> health worker shortage projected by 2030", center_body))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph("<b>98% of patients and families are left without support</b>", center_body))
    story.append(PageBreak())

    # Slide 4: Opportunity
    story.append(Paragraph("India's Digital Infrastructure: An Opportunity", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>500M+</b> WhatsApp Users", center_body))
    story.append(Paragraph("<b>1.2B</b> Mobile Subscribers", center_body))
    story.append(Paragraph("<b>22</b> Scheduled Languages", center_body))
    story.append(Paragraph("<b>Bhashini</b> AI Language Platform", center_body))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph("<b>Voice-first AI can reach anyone with a phone — no app, no internet, no literacy required</b>", center_body))
    story.append(PageBreak())

    # Slide 5: Introducing
    story.append(Paragraph("Introducing Palli Sahayak", title_style))
    story.append(Paragraph("AI Voice Helpline for Palliative Care", subtitle_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Phone Calls</b>: Dial in any language, get clinically-grounded answers", center_body))
    story.append(Paragraph("<b>WhatsApp</b>: Text or voice messages in 15+ Indian languages", center_body))
    story.append(Paragraph("<b>Web Voice</b>: Real-time voice streaming via browser", center_body))
    story.append(Paragraph("<b>Always Safe</b>: Emergency detection & professional referral", center_body))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph("<b>Answering questions at 2 AM when families need help most</b>", center_body))
    story.append(PageBreak())

    # Slide 6: Architecture
    story.append(Paragraph("System Architecture", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Input Channels:</b> Phone (Bolna.ai) | WhatsApp (Twilio) | Web Voice (Gemini Live)", center_body))
    story.append(Paragraph("↓", center_body))
    story.append(Paragraph("<b>Voice Router & Smart Query Classifier</b>", center_body))
    story.append(Paragraph("↓", center_body))
    story.append(Paragraph("<b>Hybrid RAG Pipeline</b>", ParagraphStyle('RAG', parent=center_body, textColor=PRIMARY, fontName='Helvetica-Bold')))
    story.append(Paragraph("ChromaDB Vector Search + Neo4j Knowledge Graph + GraphRAG Community Reports", center_body))
    story.append(Paragraph("↓", center_body))
    story.append(Paragraph("<b>Clinically Validated Voice/Text Response</b>", center_body))
    story.append(PageBreak())

    # Slide 7: Voice-First
    story.append(Paragraph("Why Voice-First Design?", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Literacy</b>: Voice works for everyone regardless of reading ability", center_body))
    story.append(Paragraph("<b>Technology</b>: No smartphone needed — just a basic phone call", center_body))
    story.append(Paragraph("<b>Language</b>: Speak in your mother tongue — 15+ Indian languages", center_body))
    story.append(Paragraph("<b>Geography</b>: Rural areas with voice-only connectivity", center_body))
    story.append(Paragraph("<b>Timing</b>: 24/7 availability — help at 2 AM", center_body))
    story.append(Paragraph("<b>Dignity</b>: Understood in your language, respected in your culture", center_body))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph("<b>Voice AI democratizes access to quality healthcare information</b>", center_body))
    story.append(PageBreak())

    # Slide 8: Languages
    story.append(Paragraph("15+ Indian Languages", title_style))
    story.append(Paragraph("Speaking to users in their mother tongue", subtitle_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>P0 (Priority)</b>: Hindi (हिन्दी), English, Bengali (বাংলা), Tamil (தமிழ்)", center_body))
    story.append(Paragraph("<b>P1</b>: Telugu (తెలుగు), Marathi (मराठी), Gujarati (ગુજરાતી), Kannada (ಕನ್ನಡ)", center_body))
    story.append(Paragraph("<b>P2</b>: Malayalam (മലയാളം), Punjabi (ਪੰਜਾਬੀ), Odia (ଓଡ଼ିଆ)", center_body))
    story.append(Paragraph("<b>P3</b>: Assamese (অসমীয়া) and more...", center_body))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph("<b>Powered by Bhashini — India's Government AI Language Platform</b>", center_body))
    story.append(Paragraph("Free API access for Digital Public Goods | 22 scheduled languages supported", center_body))
    story.append(PageBreak())

    # Slide 9: RAG Foundation
    story.append(Paragraph("Clinical RAG Foundation", title_style))
    story.append(Paragraph("Retrieval-Augmented Generation with Medical Guidelines", subtitle_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>1. User Question</b> → Voice or text query in any language", center_body))
    story.append(Paragraph("<b>2. Query Processing</b> → Translation + entity extraction", center_body))
    story.append(Paragraph("<b>3. Knowledge Retrieval</b> → Search clinical guidelines", center_body))
    story.append(Paragraph("<b>4. Response Generation</b> → Contextual answer with citations", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Grounded in Clinical Guidelines:</b>", center_body))
    story.append(Paragraph("WHO Palliative Care | IAHPC Essential Practices | IAPC Handbook | Pallium India Protocols", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Every response is grounded in verified medical sources</b>", center_body))
    story.append(PageBreak())

    # Slide 10: Knowledge Architecture
    story.append(Paragraph("Triple-Layer Knowledge Architecture", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>ChromaDB</b> — Vector Search", center_body))
    story.append(Paragraph("Semantic similarity with 384-dimensional embeddings", center_body))
    story.append(Spacer(1, 0.2*inch))
    story.append(Paragraph("<b>Neo4j</b> — Knowledge Graph", center_body))
    story.append(Paragraph("Structured relationships with 18 entity types", center_body))
    story.append(Spacer(1, 0.2*inch))
    story.append(Paragraph("<b>GraphRAG</b> — Community Reports", center_body))
    story.append(Paragraph("Multi-hop reasoning with Leiden clustering", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Intelligent Context Fusion</b>", center_body))
    story.append(Paragraph("FUSE: When multiple contexts are equally relevant (distance < 0.15)", center_body))
    story.append(Paragraph("FOCUS: When one context is clearly best", center_body))
    story.append(PageBreak())

    # Slide 11: Clinical Validation
    story.append(Paragraph("Clinical Validation Pipeline", title_style))
    story.append(Paragraph("Ensuring Safe, Accurate Medical Information", subtitle_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Dosage Check</b>: Validate medication doses against ranges", center_body))
    story.append(Paragraph("<b>Safety Keywords</b>: Detect emergencies and urgent situations", center_body))
    story.append(Paragraph("<b>Hallucination Check</b>: Verify response against source documents", center_body))
    story.append(Paragraph("<b>Citation Verify</b>: Ensure proper source attribution", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Expert Review Sampling:</b>", center_body))
    story.append(Paragraph("5% Normal queries | 50% High-risk topics | 100% Critical/Emergency", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Target: 95%+ Accuracy | <2% Hallucination | 90%+ Expert Agreement</b>", center_body))
    story.append(PageBreak())

    # Slide 12: Personalization
    story.append(Paragraph("Adaptive User Experience", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Automatic Role Detection</b>", center_body))
    story.append(Paragraph("Patient: Personal symptom management (\"मुझे दर्द हो रहा है\")", center_body))
    story.append(Paragraph("Caregiver: Care guidance for family member (\"my mother has...\")", center_body))
    story.append(Paragraph("Healthcare Worker: Clinical decision support (\"my patient needs...\")", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Context Memory Across Sessions</b>", center_body))
    story.append(Paragraph("Patient Condition | Current Symptoms | Medications | Allergies", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<i>\"Last time you asked about pain management...\" — Continuity of care</i>", center_body))
    story.append(PageBreak())

    # Slide 13: Analytics
    story.append(Paragraph("Real-time Analytics Dashboard", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Response Time</b>: P50: 400ms | P95: 800ms | P99: 1.2s", center_body))
    story.append(Paragraph("<b>RAG Success</b>: 89.3% retrieval rate", center_body))
    story.append(Paragraph("<b>Validation Rate</b>: 95.2% pass rate", center_body))
    story.append(Paragraph("<b>User Satisfaction</b>: 4.6/5.0 average", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Language Distribution</b>: Hindi 45% | Bengali 18% | Tamil 12% | English 10% | Others 15%", center_body))
    story.append(Paragraph("<b>Topic Distribution</b>: Pain Mgmt 32% | Symptoms 24% | Medications 18% | Caregiver 14% | End-of-Life 12%", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>System Health: HEALTHY | Alerts: 0 | Uptime: 99.9%</b>", center_body))
    story.append(PageBreak())

    # Slide 14: DPG
    story.append(Paragraph("A Digital Public Good (DPG)", title_style))
    story.append(Paragraph("Open Source | Replicable | Sustainable", subtitle_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("✓ <b>Open Source</b>: MIT License, Publicly on GitHub", center_body))
    story.append(Paragraph("✓ <b>SDG Aligned</b>: SDG 3: Good Health | SDG 10: Reduced Inequalities", center_body))
    story.append(Paragraph("✓ <b>Privacy First</b>: No personal health data, Anonymized interactions", center_body))
    story.append(Paragraph("✓ <b>Do No Harm</b>: Medical guardrails, Professional referral", center_body))
    story.append(Paragraph("✓ <b>Interoperable</b>: REST APIs, Modular architecture", center_body))
    story.append(Paragraph("✓ <b>Data Standards</b>: SNOMED, ICD-10, Medical ontologies", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Aligned with Digital Public Goods Alliance (DPGA) Standards</b>", center_body))
    story.append(PageBreak())

    # Slide 15: Global Replicability
    story.append(Paragraph("Beyond India: Global Replicability", title_style))
    story.append(Paragraph("A model for palliative care AI worldwide", subtitle_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Configuration-Based Adaptation</b> — Just change config.yaml", center_body))
    story.append(Spacer(1, 0.2*inch))
    story.append(Paragraph("<b>Potential Adaptations:</b>", center_body))
    story.append(Paragraph("Sub-Saharan Africa: Kenya, Nigeria, South Africa", center_body))
    story.append(Paragraph("Southeast Asia: Indonesia, Philippines, Vietnam", center_body))
    story.append(Paragraph("Latin America: Brazil, Mexico, Colombia", center_body))
    story.append(Paragraph("Middle East: Egypt, Jordan, Lebanon", center_body))
    story.append(Paragraph("Developed Nations: Rural US, UK NHS, Australia", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Partnership Opportunities</b>: WHO | UNICEF | Open Source Pharma | Local Health Ministries | IAHPC", center_body))
    story.append(PageBreak())

    # Slide 16: Impact Metrics
    story.append(Paragraph("Impact Metrics & Targets", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>V1 → V2 Improvements:</b>", center_body))
    story.append(Paragraph("Monthly Active Users: 1,000 → 100,000 (100x)", center_body))
    story.append(Paragraph("Languages Supported: 6 → 15+ (2.5x)", center_body))
    story.append(Paragraph("Response Accuracy: 85% → 95%+ (+10%)", center_body))
    story.append(Paragraph("Voice Latency: 2.5s → <0.5s (5x faster)", center_body))
    story.append(Paragraph("Healthcare Workers: 100 → 10,000 (100x)", center_body))
    story.append(Paragraph("Cost per Query: $0.02 → $0.005 (4x cheaper)", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>2-Year Vision:</b> 1M Users | 500K Care Decisions/Year | 10 Countries | 5 State Adoptions", center_body))
    story.append(PageBreak())

    # Slide 17: Partners
    story.append(Paragraph("Partners & Collaborations", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Clinical Partners</b>: Pallium India | Max Healthcare | CanSupport | DNipCare", center_body))
    story.append(Paragraph("<b>Academic</b>: Ashoka University | KCDH | ASAR/Nivarana", center_body))
    story.append(Paragraph("<b>Government</b>: BIRAC-DBT | Grand Challenges India | Bhashini", center_body))
    story.append(Paragraph("<b>Funders</b>: Gates Foundation | BIRAC", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<i>Mentorship: Padma Shri Dr. M.R. Rajagopal — Father of Palliative Care in India</i>", center_body))
    story.append(PageBreak())

    # Slide 18: Tech Stack
    story.append(Paragraph("Technology Stack", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Framework</b>: FastAPI • Uvicorn • Gradio Admin", center_body))
    story.append(Paragraph("<b>Voice AI</b>: Bolna.ai • Gemini Live • Groq Whisper • Edge TTS", center_body))
    story.append(Paragraph("<b>RAG Pipeline</b>: ChromaDB • Neo4j • Microsoft GraphRAG", center_body))
    story.append(Paragraph("<b>LLMs</b>: Groq (llama-3.1-8b) • Qwen3-32B • GPT-4o-mini", center_body))
    story.append(Paragraph("<b>Languages</b>: Bhashini • Deepgram • ElevenLabs", center_body))
    story.append(Paragraph("<b>Infrastructure</b>: Python 3.11 • Docker • ngrok", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>100% Open Source • github.com/inventcures/rag_gci • MIT License</b>", center_body))
    story.append(PageBreak())

    # Slide 19: Indian Context
    story.append(Paragraph("Situated in Indian Socio-Cultural Context", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Family-Centered Care</b>: Healthcare decisions are family matters. Palli Sahayak supports family caregivers.", center_body))
    story.append(Paragraph("<b>Spiritual Dimensions</b>: End-of-life care integrates spiritual practices. Respectful of diverse beliefs.", center_body))
    story.append(Paragraph("<b>Home-Based Care</b>: Most palliative care happens at home. Designed for caregivers, not just professionals.", center_body))
    story.append(Paragraph("<b>Language Diversity</b>: India's linguistic diversity is respected. 15+ languages with cultural responses.", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>India-Specific Case Scenarios Dataset</b>", center_body))
    story.append(Paragraph("Curated with Pallium India & Max Healthcare — Reflecting local clinical contexts", center_body))
    story.append(PageBreak())

    # Slide 20: Roadmap
    story.append(Paragraph("Development Roadmap", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Q1-Q2 2026: Scale & Validate</b>", center_body))
    story.append(Paragraph("GraphRAG integration | Clinical validation | 15 language rollout | Sub-500ms latency", center_body))
    story.append(Spacer(1, 0.2*inch))
    story.append(Paragraph("<b>Q3-Q4 2026: Expand & Partner</b>", center_body))
    story.append(Paragraph("100K user target | 5+ hospice integrations | Research publications | State pilot programs", center_body))
    story.append(Spacer(1, 0.2*inch))
    story.append(Paragraph("<b>2027+: Global Impact</b>", center_body))
    story.append(Paragraph("1M users target | 10 country adaptations | Policy integration | Research platform", center_body))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Current Status: V2 Development — Clinical Validation, Analytics, Personalization</b>", center_body))
    story.append(PageBreak())

    # Slide 21: Call to Action
    story.append(Paragraph("How You Can Help", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<b>Healthcare Professionals</b>: Validate clinical content | Provide feedback | Refer patients and caregivers", center_body))
    story.append(Paragraph("<b>Researchers</b>: Collaborate on effectiveness studies | Contribute to datasets | Co-author publications", center_body))
    story.append(Paragraph("<b>Technologists</b>: Contribute to open source code | Adapt for your region | Build integrations", center_body))
    story.append(Paragraph("<b>Funders & Partners</b>: Support infrastructure scaling | Fund regional adaptations | Policy advocacy", center_body))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph("<b>GitHub</b>: github.com/inventcures/rag_gci", center_body))
    story.append(Paragraph("<b>Docs</b>: deepwiki.com/inventcures/rag_gci", center_body))
    story.append(PageBreak())

    # Slide 22: Thank You
    story.append(Spacer(1, 1*inch))
    story.append(Paragraph("Thank You", ParagraphStyle('ThankYou', parent=title_style, fontSize=48)))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<i>Palli Sahayak — Compassionate AI for Palliative Care</i>", subtitle_style))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph("<b>Contact & Resources</b>", center_body))
    story.append(Paragraph("GitHub: github.com/inventcures/rag_gci", center_body))
    story.append(Paragraph("Documentation: deepwiki.com/inventcures/rag_gci", center_body))
    story.append(Paragraph("Grant: gcgh.grandchallenges.org", center_body))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph("<b>Open Source • MIT License • A Digital Public Good for India and the World</b>", center_body))

    # Build PDF
    doc.build(story)
    return pdf_path


def main():
    """Generate the presentation."""
    print("Creating Palli Sahayak presentation...")

    prs = create_presentation()

    # Save PPTX
    output_dir = os.path.dirname(os.path.abspath(__file__))
    pptx_path = os.path.join(output_dir, "Palli_Sahayak_Presentation.pptx")
    prs.save(pptx_path)
    print(f"Saved PPTX: {pptx_path}")

    # Save PDF
    print("\nGenerating PDF version...")
    pdf_path = create_pdf_presentation()
    print(f"Saved PDF: {pdf_path}")

    print(f"\nPresentation created with {len(prs.slides)} slides:")
    print("1.  Title Slide")
    print("2.  Acknowledgements (Sponsors & Partners)")
    print("3.  The Palliative Care Crisis")
    print("4.  India's Digital Opportunity")
    print("5.  Introducing Palli Sahayak")
    print("6.  System Architecture")
    print("7.  Why Voice-First Design")
    print("8.  15+ Indian Languages")
    print("9.  Clinical RAG Foundation")
    print("10. Triple-Layer Knowledge Architecture")
    print("11. Clinical Validation Pipeline")
    print("12. Adaptive User Experience")
    print("13. Real-time Analytics Dashboard")
    print("14. Digital Public Good (DPG)")
    print("15. Global Replicability")
    print("16. Impact Metrics & Targets")
    print("17. Partners & Collaborations")
    print("18. Technology Stack")
    print("19. Indian Socio-Cultural Context")
    print("20. Development Roadmap")
    print("21. Call to Action")
    print("22. Thank You")

    return pptx_path, pdf_path


if __name__ == "__main__":
    main()
