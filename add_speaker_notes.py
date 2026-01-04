#!/usr/bin/env python3
"""
Add speaker notes to Palli Sahayak presentation (PPTX).
"""

from pptx import Presentation
import os

# Speaker notes for each slide (plain text format for PPTX)
SPEAKER_NOTES = [
    # Slide 1: Title
    """[Duration: 30 sec]

<HOOK - Start with a question>

"Imagine it's 2 AM. Your mother is in pain. Cancer pain. You don't know what to do. The doctor's clinic is closed. Google gives you 50 contradicting answers.

<pause - 3 sec>

What do you do?"

<pause - 2 sec>

"This is the reality for over 10 million Indian families RIGHT NOW. Today, I'm going to show you how we're changing that."

[Stage direction: Walk to center stage, make eye contact]""",

    # Slide 2: Acknowledgements
    """[Duration: 45 sec]

"Before we dive in, I want to acknowledge the shoulders we stand on.

This work is possible because of the Bill & Melinda Gates Foundation and BIRAC through Grand Challenges India.

<pause - 2 sec>

Our PI, Dr. Anurag Agrawal - Dean of Trivedi School of Biosciences and Head of KCDH-Ashoka - whose vision drives this project.

And our incredible clinical partners: Pallium India and IAPC, who ensure everything we build actually WORKS in the real world."

[Stage direction: Nod to acknowledge, brief but sincere]""",

    # Slide 3: The Crisis
    """[Duration: 50 sec]

<Poll the audience>

"Quick show of hands - how many of you have had a family member who needed palliative care?"

<pause - 3 sec, scan the room>

"Now, let these numbers sink in:

10 MILLION Indians need palliative care. Only 1-2% receive it.

<pause - 2 sec>

That's not a gap. That's a CHASM.

We have 0.4 palliative care specialists per 100,000 people. The US has 3. The UK has 5.

But here's the thing - we can't train doctors fast enough. So what's the alternative?"

<pause - 2 sec>""",

    # Slide 4: India's Digital Opportunity
    """[Duration: 40 sec]

"Here's where India has a secret superpower.

<Can someone tell me how many active mobile connections India has?>

<pause - 3 sec>

1.2 BILLION. More phones than toilets, as they say. [slight smile]

And here's the beautiful part - voice calls work EVERYWHERE. No 4G needed. No smartphone needed. No literacy needed.

The phone is already in every hand. We just need to make it smarter.""",

    # Slide 5: Introducing Palli Sahayak
    """[Duration: 45 sec]

"So we built Palli Sahayak - which means 'Companion in Care' in Hindi.

It's NOT a chatbot. <pause - 1 sec> It's NOT a symptom checker. <pause - 1 sec>

It's a voice-first AI helpline that speaks YOUR language - literally 15+ Indian languages - available 24/7, answering the question: 'What do I do RIGHT NOW?'

<Ask a q> What makes this different from Dr. Google?

<pause - 2 sec>

Exactly - Clinical validation. Every response is grounded in palliative care guidelines, checked for safety, verified by experts. This isn't generic health info - it's ACTIONABLE guidance.""",

    # Slide 6: System Architecture
    """[Duration: 40 sec]

"Let me show you what's under the hood - but I promise to keep it simple.

[Point to diagram]

You call or WhatsApp. Your voice goes through one of three swappable AI providers - we'll see those next. Gets transcribed, processed through our RAG pipeline - that's where the clinical knowledge lives - and comes back as natural speech.

The magic? Under 1 second latency. It feels like talking to a person, not waiting for a computer.

<pause - 2 sec>

Now let me show you the REALLY cool part...""",

    # Slide 7: LIVE DEMO
    """[Duration: 30 sec intro + 10 min demo]

"Alright, enough slides. Let's see this thing in action."

<DEMO TIME - 10 minutes>

Demo sequence:
1. Start with Hindi voice call - "Meri maa ko bahut dard ho raha hai" (2 min)
2. Switch to English WhatsApp - medication question (2 min)
3. Show Hinglish interaction - code-switching (2 min)
4. Demonstrate safety guardrails - show what happens with emergency query (2 min)
5. Show admin dashboard - real-time analytics (2 min)

<pause after demo - 3 sec>

"Questions so far before we continue?""",

    # Slide 8: Voice AI Providers
    """[Duration: 50 sec]

"Here's something I'm particularly proud of - our provider-agnostic architecture.

We don't lock into one vendor. Three providers, completely swappable:

[Point to each]

Gemini Live - Google's native audio, 4 Indian languages, incredibly natural.

Bolna.ai - 5 languages including Hinglish - that beautiful mix of Hindi and English that 400 million Indians actually speak.

Retell.ai with Vobiz - real Indian phone numbers, +91, works on any phone.

<Can someone tell me why having multiple providers matters?>

<pause - 3 sec>

Exactly - redundancy, cost optimization, and no vendor lock-in. This is how you build infrastructure that LASTS.""",

    # Slide 9: Voice-First Design
    """[Duration: 35 sec]

"Why voice? <pause - 1 sec>

Because dignity means being understood in your mother tongue.

<Ask a q> How many of your parents are comfortable typing in English on a smartphone?

<pause - 2 sec>

Voice is universal. Voice is natural. Voice works at 2 AM when you're exhausted and scared.

No app download. No login. No typing. Just... talk.""",

    # Slide 10: Languages
    """[Duration: 30 sec]

"Hindi, Tamil, Telugu, Bengali, Marathi, Kannada... and yes, Hinglish.

[Smile]

Because let's be honest - most of urban India doesn't speak 'pure' Hindi or 'pure' English. We code-switch constantly. 'Mujhe morning mein pain hota hai' - that's real speech.

Our system understands that.""",

    # Slide 11: Clinical RAG Foundation
    """[Duration: 40 sec]

"Now, the brain of the system - our Clinical RAG.

RAG stands for Retrieval Augmented Generation. Fancy words, simple concept:

Instead of the AI making things up, it RETRIEVES real information from curated palliative care knowledge, THEN generates a response.

<pause - 2 sec>

We're talking WHO guidelines, IAPC protocols, Pallium India resources - over 200 verified documents. Not Reddit posts. Not random blogs.""",

    # Slide 12: Triple-Layer Architecture
    """[Duration: 45 sec]

"Here's where it gets interesting. We don't just do simple search.

Three layers: [count on fingers]

ONE - Vector search. Semantic similarity. 'Pain' matches 'discomfort' matches 'it hurts'.

TWO - Knowledge Graph. Relationships. 'Morphine treats pain' but 'Morphine requires prescription' and 'Morphine has side effects'.

THREE - GraphRAG. Community detection. Understanding CONTEXT across multiple documents.

<pause - 2 sec>

Then we FUSE these layers intelligently. That's the secret sauce.""",

    # Slide 13: Clinical Validation
    """[Duration: 40 sec]

<Poll the audience>

"Who here trusts AI to give medical advice without human oversight?"

<pause - 3 sec, expect few hands>

"Good. Neither do we.

Every response goes through: dosage verification, drug interaction checks, emergency detection, and periodic expert sampling.

The AI doesn't replace doctors. It EXTENDS their reach to 2 AM, to remote villages, to overwhelmed caregivers.""",

    # Slide 14: Personalization
    """[Duration: 30 sec]

"The system learns. Not in a creepy way - in a helpful way.

It remembers: this patient is on morphine. This caregiver prefers Hindi. This family has asked about anxiety before.

Context persistence. So you don't repeat your story every single time.""",

    # Slide 15: Analytics
    """[Duration: 30 sec]

"And for program managers - real-time visibility.

What are people asking about? Where are the pain points - literally? Which regions need more support?

Data-driven palliative care. Novel concept, right?"

[Brief smile]""",

    # Slide 16: DPG
    """[Duration: 50 sec]

"Now, here's the BIGGER picture.

<Ask a q> What is a Digital Public Good?

<pause - 3 sec>

Open source. Open standards. Designed for replication.

Palli Sahayak isn't just a project - it's INFRASTRUCTURE. Like Aadhaar. Like UPI. But for healthcare.

Any government, any NGO, any healthcare system can take this, adapt it, deploy it. The code is on GitHub. The docs are public. The knowledge is shared.""",

    # Slide 17: Global Replicability
    """[Duration: 45 sec]

"And this isn't just for India.

<Can someone tell me which developed country has perfect palliative care coverage?>

<pause - 3 sec>

None. The US has 1 million people who need hospice but don't get it. UK has rural access issues. Japan has an aging crisis.

This architecture - voice-first, multilingual, clinically validated - works ANYWHERE. Replace Hindi with Spanish, replace IAPC guidelines with NHPCO guidelines. The bones are the same.""",

    # Slide 18: Impact Metrics
    """[Duration: 35 sec]

"What does success look like?

Pilot: 1,000 families. Measure: symptom resolution time, caregiver stress, unnecessary ER visits.

Scale: 100,000 families in 24 months.

Moon shot: Integrate with National Health Mission. Make this the palliative care backbone for Ayushman Bharat.

<pause - 2 sec>

Ambitious? Yes. Impossible? Watch us.""",

    # Slide 19: Partners
    """[Duration: 30 sec]

"We don't build alone.

Pallium India - clinical expertise, training data, real-world validation.

IAPC - the Indian Association of Palliative Care - guidelines, network, credibility.

This is what responsible AI development looks like - technologists and clinicians, together.""",

    # Slide 20: Tech Stack
    """[Duration: 30 sec]

"For the techies in the room - yes, it's Python. Yes, it's FastAPI. Yes, it's open source.

ChromaDB for vectors, Neo4j for knowledge graphs, multiple LLM backends.

The stack is boring on purpose. Boring means maintainable. Boring means any developer can contribute.""",

    # Slide 21: Indian Context
    """[Duration: 40 sec]

"But technology is only half the story.

We designed for INDIA. Joint families where decisions are collective. Stigma around morphine. Regional dietary restrictions during illness. Spiritual needs alongside medical ones.

<Ask a q> How many medical AI systems consider that a Hindu family might need guidance on rituals during end-of-life?

<pause - 2 sec>

Ours does.""",

    # Slide 22: Roadmap
    """[Duration: 30 sec]

"Where are we going?

Next: Pilot in 3 states. Then: Integration with telemedicine platforms. Eventually: A national helpline number that anyone can call.

The roadmap is aggressive because the need is urgent.""",

    # Slide 23: Call to Action
    """[Duration: 45 sec]

<CLOSING - Bring energy back up>

"So what can YOU do?

If you're a clinician - help us validate. Your expertise makes this safe.

If you're a researcher - collaborate. There's a dozen papers waiting to be written.

If you're a funder - invest. Not in a startup, in INFRASTRUCTURE.

If you're anyone else - spread the word. Someone you know needs this."

<pause - 2 sec>""",

    # Slide 24: Thank You
    """[Duration: 30 sec + 5 min Q&A]

<CLOSING HOOK - Circle back to opening>

"Remember that 2 AM scenario I started with?

<pause - 2 sec>

Now imagine: the caregiver picks up the phone, speaks in Hindi, gets clear guidance on managing the pain, knows when to call the doctor, and gets through the night.

THAT'S what we're building.

Thank you. I'll take questions."

<Q&A - 5 minutes>""",
]

def add_notes_to_pptx(pptx_path, output_path):
    """Add speaker notes to each slide in the presentation."""
    prs = Presentation(pptx_path)

    for i, slide in enumerate(prs.slides):
        if i < len(SPEAKER_NOTES):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = SPEAKER_NOTES[i]
            print(f"Added notes to slide {i+1}")

    prs.save(output_path)
    print(f"\nSaved presentation with notes: {output_path}")

if __name__ == "__main__":
    input_path = "/Users/tp53/Documents/tp53_AA/llms4palliative_gci/06JAN2026___WORKSHOP/demo+my_ppt/v3_Palli_Sahayak_Presentation.pptx"
    output_path = "/Users/tp53/Documents/tp53_AA/llms4palliative_gci/06JAN2026___WORKSHOP/demo+my_ppt/v3_Palli_Sahayak_Presentation.pptx"
    add_notes_to_pptx(input_path, output_path)
