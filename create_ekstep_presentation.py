#!/usr/bin/env python3
"""
Palli Sahayak - EkStep Voice AI Event Presentation
==================================================

A demo-first showcase presentation for India's best production-grade 
Voice AI deployment at the EkStep Foundation event.

Event: Voice AI - Making the Best Work for India
Venue: The Ritz-Carlton, Bengaluru
Date: January 28, 2026

Author: Palli Sahayak Team
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme based on the event poster (brown/earthy tones)
PRIMARY_COLOR = RgbColor(0x6B, 0x4E, 0x3D)      # Dark brown
SECONDARY_COLOR = RgbColor(0xD4, 0xA5, 0x7D)    # Light brown/tan
ACCENT_COLOR = RgbColor(0xF5, 0xC1, 0x6C)       # Gold/yellow accent
WHITE = RgbColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RgbColor(0x3E, 0x2C, 0x22)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_COLOR
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, has_demo=False, demo_text=""):
    """Add a content slide with bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(7.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"●  {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(16)
        p.level = 0
    
    # Demo box if needed
    if has_demo:
        demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.5), Inches(4.5), Inches(5.5))
        demo_box.fill.solid()
        demo_box.fill.fore_color.rgb = SECONDARY_COLOR
        demo_box.line.color.rgb = ACCENT_COLOR
        demo_box.line.width = Pt(3)
        
        demo_textbox = slide.shapes.add_textbox(Inches(8.7), Inches(1.7), Inches(4.1), Inches(5.1))
        tf = demo_textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "🎬 LIVE DEMO"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        p = tf.add_paragraph()
        p.text = demo_text
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(20)
    
    return slide

def add_image_slide(prs, title, description):
    """Add a slide with large text for visual impact"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = SECONDARY_COLOR
    bg.line.fill.background()
    
    # Large title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(1.5))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_architecture_slide(prs):
    """Add architecture diagram slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Technical Architecture"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Architecture boxes - simplified representation
    y_start = Inches(1.8)
    box_height = Inches(1.0)
    box_width = Inches(3.5)
    
    # Row 1: Voice Providers
    providers = [
        ("🎙️ Gemini Live\n(Web Voice)", RgbColor(0x90, 0xEE, 0x90)),
        ("📞 Bolna.ai\n(Phone Calls)", RgbColor(0x87, 0xCE, 0xEB)),
        ("📱 Retell + Vobiz\n(PSTN +91)", RgbColor(0xDD, 0xA0, 0xDD)),
    ]
    
    for i, (text, color) in enumerate(providers):
        x = Inches(0.8 + i * 4.0)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_start, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)
        
        tb = slide.shapes.add_textbox(x + Inches(0.1), y_start + Inches(0.2), box_width - Inches(0.2), box_height - Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER
    
    # Arrow down
    arrow_y = y_start + box_height + Inches(0.2)
    
    # Voice Router
    router_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), arrow_y, Inches(3.5), Inches(0.8))
    router_box.fill.solid()
    router_box.fill.fore_color.rgb = ACCENT_COLOR
    router_box.line.color.rgb = PRIMARY_COLOR
    router_box.line.width = Pt(3)
    
    tb = slide.shapes.add_textbox(Inches(5.0), arrow_y + Inches(0.2), Inches(3.3), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "🔀 Voice Router + Safety Layer"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # RAG Pipeline
    rag_y = arrow_y + Inches(1.2)
    rag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), rag_y, Inches(8.3), Inches(1.2))
    rag_box.fill.solid()
    rag_box.fill.fore_color.rgb = RgbColor(0xF0, 0xE6, 0x8C)
    rag_box.line.color.rgb = PRIMARY_COLOR
    rag_box.line.width = Pt(2)
    
    tb = slide.shapes.add_textbox(Inches(2.7), rag_y + Inches(0.15), Inches(7.9), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏥 RAG Pipeline + Knowledge Graph + Longitudinal Memory"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # Bottom: Data sources
    data_y = rag_y + Inches(1.5)
    data_sources = [
        ("📚 Medical Corpus", Inches(0.8)),
        ("🧠 Knowledge Graph", Inches(4.5)),
        ("💊 Patient History", Inches(8.2)),
    ]
    
    for text, x in data_sources:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, data_y, Inches(3.8), Inches(0.7))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(1)
        
        tb = slide.shapes.add_textbox(x + Inches(0.1), data_y + Inches(0.15), Inches(3.6), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_stats_slide(prs):
    """Add impact stats slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_COLOR
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Impact & Scale"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Stats boxes
    stats = [
        ("5", "Indian Languages", "hi, bn, ta, gu, en"),
        ("3", "Voice AI Providers", "Gemini, Bolna, Retell"),
        ("24/7", "Availability", "Always-on helpline"),
        ("<2s", "Response Time", "Real-time voice AI"),
    ]
    
    positions = [
        (Inches(1), Inches(2)),
        (Inches(7), Inches(2)),
        (Inches(1), Inches(4.5)),
        (Inches(7), Inches(4.5)),
    ]
    
    for i, (number, label, subtext) in enumerate(stats):
        x, y = positions[i]
        
        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5), Inches(2))
        box.fill.solid()
        box.fill.fore_color.rgb = SECONDARY_COLOR
        box.line.color.rgb = ACCENT_COLOR
        box.line.width = Pt(2)
        
        # Number
        tb = slide.shapes.add_textbox(x, y + Inches(0.2), Inches(5), Inches(0.8))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        tb = slide.shapes.add_textbox(x, y + Inches(1.0), Inches(5), Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER
        
        # Subtext
        tb = slide.shapes.add_textbox(x, y + Inches(1.5), Inches(5), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = subtext
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER
    
    return slide

# ==================== BUILD PRESENTATION ====================

# Slide 1: Title
add_title_slide(prs, 
    "Palli Sahayak",
    "Voice AI for Palliative Care | Making Healthcare Accessible to Every Indian")

# Slide 2: The Problem
add_image_slide(prs,
    "The Challenge",
    "10M+ Indians need palliative care | 1% have access | Language barriers | 24/7 support needed")

# Slide 3: Our Solution
add_content_slide(prs,
    "Palli Sahayak: India's First Voice AI Palliative Care Helpline",
    [
        "Multi-modal: WhatsApp, Phone, Web Voice",
        "5 Indian languages: Hindi, Bengali, Tamil, Gujarati, English",
        "3 Voice AI providers for redundancy & coverage",
        "RAG-powered medical knowledge base",
        "Emergency detection & human handoff",
        "Medication reminders with voice calls",
    ],
    has_demo=True,
    demo_text="1. Live voice call demo\n2. Hindi conversation\n3. Emergency escalation\n4. Medication reminder")

# Slide 4: Technical Architecture
add_architecture_slide(prs)

# Slide 5: Safety First (NEW FEATURES)
add_content_slide(prs,
    "Safety-First AI: 5 Production-Grade Enhancements",
    [
        "🚨 Emergency Detection: Auto-escalation in 5 languages",
        "🔬 Evidence Badges: Confidence scores on every response",
        "📏 Smart Responses: Length-optimized for comprehension",
        "👤 Human Handoff: One-click transfer to caregivers",
        "💊 Voice Reminders: Automated medication call-backs",
    ],
    has_demo=True,
    demo_text="● Emergency keyword detection\n● Voice call reminder\n● Patient confirmation")

# Slide 6: Multi-Provider Voice AI
add_content_slide(prs,
    "Voice AI Router: Best of Breed Architecture",
    [
        "🎙️ Gemini Live: Web-native voice, real-time streaming",
        "📞 Bolna.ai: Phone calls, Indian accent TTS",
        "📱 Retell + Vobiz: PSTN connectivity, +91 numbers",
        "🔄 Smart routing: Automatic failover between providers",
        "🛡️ Unified safety layer across all channels",
    ])

# Slide 7: Impact & Scale
add_stats_slide(prs)

# Slide 8: Demo Flow
add_content_slide(prs,
    "Live Demo: Patient Journey",
    [
        "Patient calls via phone or WhatsApp voice",
        "AI detects language automatically (Hindi)",
        "Patient: 'माँ को दर्द है, क्या करूं?'",
        "AI queries RAG + Knowledge Graph",
        "Voice response with evidence badge",
        "Emergency keywords trigger instant escalation",
    ],
    has_demo=True,
    demo_text="⏯️ Starting live demo...\n\n📞 Call initiated\n🗣️ Hindi voice input\n🔍 Knowledge retrieval\n📢 Voice response")

# Slide 9: Why This Matters
add_image_slide(prs,
    "Democratizing Healthcare",
    "No smartphone needed | Works on ₹500 feature phones | Speak in your language | Available at 2 AM")

# Slide 10: Future Roadmap
add_content_slide(prs,
    "Roadmap: What's Next",
    [
        "🏥 FHIR integration with hospital EHRs",
        "📊 Predictive analytics for crisis prevention",
        "👨‍⚕️ Caregiver coordination platform",
        "🎓 Training module for community health workers",
        "🌐 Expansion: Kannada, Malayalam, Telugu",
        "🤝 Partnerships with 100+ hospice centers",
    ])

# Slide 11: Team & Recognition
add_content_slide(prs,
    "Team & Partners",
    [
        "Built with ❤️ for India's palliative care community",
        "Open source: github.com/inventcures/rag_gci",
        "Powered by: Groq, Google Gemini, Deepgram, ElevenLabs",
        "Inspired by: WHO Palliative Care Guidelines",
        "Supported by: Medical advisors from AIIMS, Tata Memorial",
        "DeepWiki: deepwiki.com/inventcures/rag_gci",
    ])

# Slide 12: Thank You
add_title_slide(prs,
    "धन्यवाद | Thank You",
    "Palli Sahayak - Your Companion in Care")

# Save presentation
output_path = "Palli_Sahayak_EkStep_Voice_AI_Presentation.pptx"
prs.save(output_path)

print(f"✅ Presentation created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"🎨 Color scheme: Earthy tones matching EkStep branding")
print(f"🎯 Focus: Demo-first showcase format")
