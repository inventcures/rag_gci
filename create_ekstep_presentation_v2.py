#!/usr/bin/env python3
"""
Palli Sahayak - EkStep Voice AI Event Presentation (ENHANCED V2)
================================================================

Incorporating:
- Website data from https://inventcures.github.io/palli-sahayak/
- GCI Grant win details (Dr Anurag Agrawal - PI, Ashish Makani - Co-I)
- Funders: GF India & BIRAC-DBT
- Co-design partners: Max Healthcare & Pallium India
- Video: Marathi demo with Dr Sachin & Dr Prakash of Cipla Foundation
- Today's safety features
- Ashoka & KCDHA affiliation

Event: Voice AI - Making the Best Work for India
Venue: The Ritz-Carlton, Bengaluru
Date: January 28, 2026

"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR
import os

# Video file path
VIDEO_PATH = "/Users/tp53/Downloads/WhatsApp Video 2026-01-05 at 19.41.26.mp4"

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# EkStep Event Color Scheme (from poster)
PRIMARY_BROWN = RGBColor(0x6B, 0x4E, 0x3D)      # Dark brown
SECONDARY_TAN = RGBColor(0xD4, 0xA5, 0x7D)      # Light brown/tan
ACCENT_GOLD = RGBColor(0xF5, 0xC1, 0x6C)        # Gold/yellow
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x3E, 0x2C, 0x22)
LIGHT_BG = RGBColor(0xF5, 0xF0, 0xEB)

def add_title_slide(prs, title, subtitle, tagline=""):
    """Add a title slide with branding"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_BROWN
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    if tagline:
        tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(12.333), Inches(0.6))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = tagline
        p.font.size = Pt(18)
        p.font.color.rgb = SECONDARY_TAN
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, highlight_text="", note=""):
    """Add a content slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BROWN
    header.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20) if not bullet.startswith("🏆") else Pt(22)
        p.font.color.rgb = DARK_TEXT if not bullet.startswith("🏆") else PRIMARY_BROWN
        p.font.bold = bullet.startswith("🏆") or bullet.startswith("⭐")
        p.space_after = Pt(14)
        p.level = 0
    
    # Highlight box if provided
    if highlight_text:
        h_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.8), Inches(11.9), Inches(1.3))
        h_box.fill.solid()
        h_box.fill.fore_color.rgb = ACCENT_GOLD
        h_box.line.color.rgb = PRIMARY_BROWN
        h_box.line.width = Pt(2)
        
        h_text = slide.shapes.add_textbox(Inches(0.9), Inches(6.0), Inches(11.5), Inches(1.0))
        tf = h_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = highlight_text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER
    
    # Note at bottom
    if note:
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.4))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(10)
        p.font.color.rgb = SECONDARY_TAN
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_full_image_slide(prs, title, subtitle, bg_color=SECONDARY_TAN):
    """Add a big text slide for impact"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(1.5))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_architecture_slide(prs):
    """Architecture diagram slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BROWN
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Multi-Provider Voice AI Architecture"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Architecture boxes - 3 providers
    y_pos = Inches(1.8)
    providers = [
        ("🎙️ Gemini Live\nWeb Voice", RGBColor(0x90, 0xEE, 0x90)),
        ("📞 Bolna.ai\nPhone Calls", RGBColor(0x87, 0xCE, 0xEB)),
        ("📱 Retell+Vobiz\nPSTN +91", RGBColor(0xDD, 0xA0, 0xDD)),
    ]
    
    for i, (text, color) in enumerate(providers):
        x = Inches(0.8 + i * 4.2)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, Inches(3.6), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = PRIMARY_BROWN
        box.line.width = Pt(2)
        
        tb = slide.shapes.add_textbox(x + Inches(0.1), y_pos + Inches(0.15), Inches(3.4), Inches(0.7))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER
    
    # Voice Router
    router_y = y_pos + Inches(1.4)
    router = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), router_y, Inches(3.5), Inches(0.7))
    router.fill.solid()
    router.fill.fore_color.rgb = ACCENT_GOLD
    router.line.color.rgb = PRIMARY_BROWN
    router.line.width = Pt(2)
    
    tb = slide.shapes.add_textbox(Inches(5.0), router_y + Inches(0.15), Inches(3.3), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "🔀 Voice Router + Safety Layer"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # RAG Pipeline
    rag_y = router_y + Inches(1.0)
    rag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), rag_y, Inches(9.3), Inches(1.0))
    rag.fill.solid()
    rag.fill.fore_color.rgb = RGBColor(0xF0, 0xE6, 0x8C)
    rag.line.color.rgb = PRIMARY_BROWN
    rag.line.width = Pt(2)
    
    tb = slide.shapes.add_textbox(Inches(2.2), rag_y + Inches(0.15), Inches(8.9), Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏥 RAG Pipeline + Knowledge Graph + Longitudinal Memory"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # Bottom layer - Data sources
    data_y = rag_y + Inches(1.3)
    sources = [
        ("📚 Indian Palliative Care\nCase Vignettes", Inches(0.8)),
        ("🧠 Neo4j\nKnowledge Graph", Inches(4.8)),
        ("💊 Patient Context\nMemory (1-5 years)", Inches(8.8)),
    ]
    
    for text, x in sources:
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, data_y, Inches(3.8), Inches(0.8))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PRIMARY_BROWN
        box.line.width = Pt(1)
        
        tb = slide.shapes.add_textbox(x + Inches(0.1), data_y + Inches(0.1), Inches(3.6), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_grant_slide(prs):
    """Grant and funding slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🏆 Grand Challenges India Grant Awardee"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    p.alignment = PP_ALIGN.CENTER
    
    # Grant details box
    grant_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.6), Inches(10.333), Inches(2.2))
    grant_box.fill.solid()
    grant_box.fill.fore_color.rgb = WHITE
    grant_box.line.color.rgb = ACCENT_GOLD
    grant_box.line.width = Pt(3)
    
    tb = slide.shapes.add_textbox(Inches(1.7), Inches(1.8), Inches(9.933), Inches(1.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Grand Challenges Global Health"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    
    p = tf.add_paragraph()
    p.text = "Clinical Decision Support Tool Using LLMs"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    p.space_before = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "Awarded: November 2024"
    p.font.size = Pt(16)
    p.font.color.rgb = SECONDARY_TAN
    p.space_before = Pt(8)
    
    # Investigators
    inv_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(1.2))
    tf = inv_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Principal Investigator: Dr. Anurag Agrawal"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "Co-Investigator: Ashish Makani"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(6)
    
    # Institution
    inst_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(12.333), Inches(0.8))
    tf = inst_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ashoka University & KCDHA (Karnataka)"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # Funders
    fund_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.8))
    tf = fund_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Funded by: Gates Foundation (India) & BIRAC-DBT"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_partner_slide(prs):
    """Clinical partners slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BROWN
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Co-Designed with Clinical Experts"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Max Healthcare
    max_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.6), Inches(5.2), Inches(2.5))
    max_box.fill.solid()
    max_box.fill.fore_color.rgb = WHITE
    max_box.line.color.rgb = PRIMARY_BROWN
    max_box.line.width = Pt(2)
    
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(4.8), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🏥 Max Healthcare"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    
    p = tf.add_paragraph()
    p.text = "New Delhi"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    
    p = tf.add_paragraph()
    p.text = "Tertiary Care Expertise"
    p.font.size = Pt(16)
    p.font.color.rgb = SECONDARY_TAN
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Case vignettes curation\n• Clinical validation\n• Quality assurance"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    
    # Pallium India
    pallium_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(1.6), Inches(5.2), Inches(2.5))
    pallium_box.fill.solid()
    pallium_box.fill.fore_color.rgb = WHITE
    pallium_box.line.color.rgb = PRIMARY_BROWN
    pallium_box.line.width = Pt(2)
    
    tb = slide.shapes.add_textbox(Inches(7.3), Inches(1.8), Inches(4.8), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "💚 Pallium India"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    
    p = tf.add_paragraph()
    p.text = "Kerala"
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT
    
    p = tf.add_paragraph()
    p.text = "Community Care Pioneers"
    p.font.size = Pt(16)
    p.font.color.rgb = SECONDARY_TAN
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Community health worker training\n• Grassroots implementation\n• Rural care protocols"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    
    # Knowledge base note
    note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.4), Inches(11.333), Inches(1.0))
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = ACCENT_GOLD
    note_box.line.color.rgb = PRIMARY_BROWN
    note_box.line.width = Pt(1)
    
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(4.6), Inches(10.933), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "📚 Knowledge Base: Indian Palliative Care Case Vignettes curated by Max Healthcare & Pallium India"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # Video reference (will be on next slide)
    video_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.7), Inches(12.333), Inches(1.5))
    tf = video_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🎬 See next slide for live demo video"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_video_slide(prs):
    """Video demo slide with embedded Marathi interaction"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_TEXT
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BROWN
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎬 Live Demo: Marathi Voice Interaction with Palliative Care Doctors"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Try to embed video if file exists
    video_embedded = False
    if os.path.exists(VIDEO_PATH):
        try:
            # Add video to slide (centered, large size)
            movie = slide.shapes.add_movie(
                VIDEO_PATH,
                Inches(1.5), Inches(1.3), Inches(10.333), Inches(5.5),
                mime_type='video/mp4'
            )
            video_embedded = True
        except Exception as e:
            print(f"Note: Could not embed video: {e}")
            video_embedded = False
    
    # If video not embedded, show placeholder with info
    if not video_embedded:
        # Placeholder box
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(1.5), Inches(1.3), Inches(10.333), Inches(5.5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
        placeholder.line.color.rgb = ACCENT_GOLD
        placeholder.line.width = Pt(3)
        
        # Play icon placeholder
        play_btn = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5.8), Inches(3.3), Inches(1.7), Inches(1.7)
        )
        play_btn.fill.solid()
        play_btn.fill.fore_color.rgb = ACCENT_GOLD
        play_btn.line.color.rgb = WHITE
        play_btn.line.width = Pt(3)
        
        # Play symbol
        play_text = slide.shapes.add_textbox(Inches(5.95), Inches(3.55), Inches(1.4), Inches(1.2))
        tf = play_text.text_frame
        p = tf.paragraphs[0]
        p.text = "▶"
        p.font.size = Pt(48)
        p.font.color.rgb = PRIMARY_BROWN
        p.alignment = PP_ALIGN.CENTER
    
    # Caption box at bottom
    caption_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.95), Inches(12.333), Inches(0.45)
    )
    caption_box.fill.solid()
    caption_box.fill.fore_color.rgb = PRIMARY_BROWN
    caption_box.line.fill.background()
    
    # Caption text
    caption = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(11.933), Inches(0.4))
    tf = caption.text_frame
    p = tf.paragraphs[0]
    p.text = "Dr Sachin & Dr Prakash, Cipla Foundation • Demonstrating natural Marathi conversation with AI"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Info note (if video embedded)
    if video_embedded:
        note = slide.shapes.add_textbox(Inches(0.5), Inches(6.65), Inches(12.333), Inches(0.25))
        tf = note.text_frame
        p = tf.paragraphs[0]
        p.text = "Click video to play • Full interaction with palliative care doctors"
        p.font.size = Pt(10)
        p.font.color.rgb = SECONDARY_TAN
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_demo_video_slide(prs):
    """Slide showcasing the system demo video (screen recording style)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_TEXT
    bg.line.fill.background()
    
    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BROWN
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎬 Live Demo: Actual Code Running in Terminal"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Try to embed the terminal demo video (shows actual code running)
    demo_video_paths = ["palli_sahayak_terminal_demo.mp4", "palli_sahayak_system_demo.mp4", "palli_sahayak_features_demo.mp4"]
    video_embedded = False
    used_path = None
    
    for demo_video_path in demo_video_paths:
        if os.path.exists(demo_video_path):
            try:
                movie = slide.shapes.add_movie(
                    demo_video_path,
                    Inches(1.5), Inches(1.3), Inches(10.333), Inches(5.5),
                    mime_type='video/mp4'
                )
                video_embedded = True
                used_path = demo_video_path
                break
            except Exception as e:
                print(f"Note: Could not embed {demo_video_path}: {e}")
                continue
    
    # If video not embedded, show placeholder
    if not video_embedded:
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.5), Inches(1.3), Inches(10.333), Inches(5.5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
        placeholder.line.color.rgb = ACCENT_GOLD
        placeholder.line.width = Pt(3)
        
        # Play icon
        play_btn = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5.8), Inches(3.3), Inches(1.7), Inches(1.7)
        )
        play_btn.fill.solid()
        play_btn.fill.fore_color.rgb = ACCENT_GOLD
        play_btn.line.color.rgb = WHITE
        play_btn.line.width = Pt(3)
        
        play_text = slide.shapes.add_textbox(Inches(5.95), Inches(3.55), Inches(1.4), Inches(1.2))
        tf = play_text.text_frame
        p = tf.paragraphs[0]
        p.text = "▶"
        p.font.size = Pt(48)
        p.font.color.rgb = PRIMARY_BROWN
        p.alignment = PP_ALIGN.CENTER
    
    # Caption at bottom
    caption_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.95), Inches(12.333), Inches(0.45)
    )
    caption_box.fill.solid()
    caption_box.fill.fore_color.rgb = PRIMARY_BROWN
    caption_box.line.fill.background()
    
    caption = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(11.933), Inches(0.4))
    tf = caption.text_frame
    p = tf.paragraphs[0]
    p.text = "Real unit tests: Oncology/COPD patients • Emergency escalation • Evidence badges with RAG citations"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_safety_slide(prs):
    """Safety features slide - highlighting today's implementation"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE7)  # Light yellow bg
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "⭐ NEW: 5 Production-Grade Safety Features"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Implemented Today - January 28, 2026"
    p.font.size = Pt(18)
    p.font.color.rgb = SECONDARY_TAN
    p.alignment = PP_ALIGN.CENTER
    
    # 5 features in boxes
    features = [
        ("🚨", "Emergency Detection", "5 languages | Auto-escalation | Caregiver alerts", Inches(0.5), Inches(1.8)),
        ("🔬", "Evidence Badges", "Confidence scores | Source quality | Physician consult warnings", Inches(6.7), Inches(1.8)),
        ("📏", "Smart Responses", "User-adaptive length | Voice-optimized | 3 comprehension levels", Inches(0.5), Inches(3.8)),
        ("👤", "Human Handoff", "Warm transfer | Request tracking | Priority queue", Inches(6.7), Inches(3.8)),
        ("💊", "Voice Reminders", "Automated calls | DTMF confirmation | Adherence tracking", Inches(3.6), Inches(5.8)),
    ]
    
    for emoji, title, desc, x, y in features:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6.0), Inches(1.7))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PRIMARY_BROWN
        box.line.width = Pt(2)
        
        # Emoji
        emoji_tb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), Inches(0.6), Inches(0.6))
        tf = emoji_tb.text_frame
        p = tf.paragraphs[0]
        p.text = emoji
        p.font.size = Pt(28)
        
        # Title
        title_tb = slide.shapes.add_textbox(x + Inches(0.8), y + Inches(0.15), Inches(5.0), Inches(0.5))
        tf = title_tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BROWN
        
        # Description
        desc_tb = slide.shapes.add_textbox(x + Inches(0.8), y + Inches(0.7), Inches(5.0), Inches(0.9))
        tf = desc_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_TEXT
    
    return slide

def add_demo_slide(prs):
    """Demo flow slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BROWN
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎬 Live Demo: Patient Journey"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Demo steps
    steps = [
        ("1", "Patient calls toll-free number", "Voice input in Hindi/Marathi"),
        ("2", "AI detects language automatically", "Natural conversation begins"),
        ("3", "Patient: 'माँ को दर्द है, क्या करूं?'", "Symptom description"),
        ("4", "System queries RAG + Knowledge Graph", "Retrieves relevant cases"),
        ("5", "Voice response with evidence badge", "Confidence score displayed"),
        ("6", "Emergency? → Instant escalation", "Safety check triggers"),
    ]
    
    y_pos = Inches(1.5)
    for num, step, detail in steps:
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y_pos, Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_GOLD
        circle.line.color.rgb = PRIMARY_BROWN
        
        tb = slide.shapes.add_textbox(Inches(0.55), y_pos + Inches(0.08), Inches(0.4), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BROWN
        p.alignment = PP_ALIGN.CENTER
        
        # Step text
        tb = slide.shapes.add_textbox(Inches(1.2), y_pos, Inches(5.5), Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        
        # Detail
        tb = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.35), Inches(5.5), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = detail
        p.font.size = Pt(14)
        p.font.color.rgb = SECONDARY_TAN
        
        y_pos += Inches(0.9)
    
    # Demo highlight box
    demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.3))
    demo_box.fill.solid()
    demo_box.fill.fore_color.rgb = ACCENT_GOLD
    demo_box.line.color.rgb = PRIMARY_BROWN
    demo_box.line.width = Pt(3)
    
    tb = slide.shapes.add_textbox(Inches(7.4), Inches(1.7), Inches(5.1), Inches(4.9))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🎥 DEMO HIGHLIGHTS"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    
    p = tf.add_paragraph()
    p.text = "\n• Voice call in Marathi\n• Doctor testing the system\n• Natural conversation flow\n• Real-time RAG retrieval\n• Safety feature activation"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.space_before = Pt(12)
    
    return slide

def add_impact_slide(prs):
    """Impact statistics slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_BROWN
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Impact & Scale"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Stats
    stats = [
        ("7+", "Indian Languages", "Hindi, Bengali, Tamil, Marathi, Punjabi, Malayalam, English"),
        ("3", "Voice AI Providers", "Gemini Live, Bolna.ai, Retell+Vobiz"),
        ("10M+", "Patients in Need", "Only 1-2% currently have access"),
        ("24/7", "Availability", "Always-on helpline"),
    ]
    
    positions = [
        (Inches(0.8), Inches(1.8)),
        (Inches(6.8), Inches(1.8)),
        (Inches(0.8), Inches(4.3)),
        (Inches(6.8), Inches(4.3)),
    ]
    
    for i, (number, label, subtext) in enumerate(stats):
        x, y = positions[i]
        
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = SECONDARY_TAN
        box.line.color.rgb = ACCENT_GOLD
        box.line.width = Pt(2)
        
        # Number
        tb = slide.shapes.add_textbox(x, y + Inches(0.15), Inches(5.7), Inches(0.9))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BROWN
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        tb = slide.shapes.add_textbox(x, y + Inches(1.05), Inches(5.7), Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER
        
        # Subtext
        tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.55), Inches(5.3), Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtext
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_open_source_slide(prs):
    """Open source and DPG slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🌍 A Digital Public Good"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    p.alignment = PP_ALIGN.CENTER
    
    # Main content
    content = [
        "Open Source: github.com/inventcures/rag_gci",
        "Documentation: deepwiki.com/inventcures/rag_gci",
        "Licensed for free use, modification, and distribution",
        "Adaptable for any country, language, or medical domain",
    ]
    
    tb = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.333), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"●  {line}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(10)
    
    # SDG Box
    sdg_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(3.9), Inches(9.333), Inches(1.3))
    sdg_box.fill.solid()
    sdg_box.fill.fore_color.rgb = ACCENT_GOLD
    sdg_box.line.color.rgb = PRIMARY_BROWN
    sdg_box.line.width = Pt(2)
    
    tb = slide.shapes.add_textbox(Inches(2.2), Inches(4.1), Inches(8.933), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Contributing to UN Sustainable Development Goals:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "SDG 3: Good Health & Well-being | SDG 10: Reduced Inequalities"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    
    # Team note
    team_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1.8))
    tf = team_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Built with ❤️ by the Palli Sahayak Team"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BROWN
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nPowered by: Groq, Google Gemini, Deepgram, ElevenLabs, ChromaDB, Neo4j"
    p.font.size = Pt(14)
    p.font.color.rgb = SECONDARY_TAN
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_thank_you_slide(prs):
    """Final thank you slide with logos reference"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_BROWN
    bg.line.fill.background()
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "धन्यवाद"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # English
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(48)
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tag_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12.333), Inches(0.8))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Palli Sahayak: Your Companion in Care"
    p.font.size = Pt(22)
    p.font.color.rgb = SECONDARY_TAN
    p.alignment = PP_ALIGN.CENTER
    
    # Contact
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.333), Inches(0.8))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "inventcures.github.io/palli-sahayak"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ==================== BUILD PRESENTATION ====================

print("Building Palli Sahayak EkStep Presentation v2...")

# Slide 1: Title
add_title_slide(prs, 
    "पल्ली सहायक  •  Palli Sahayak",
    "Voice AI for Palliative Care",
    "Democratizing Healthcare Access Across India")

# Slide 2: The Problem
add_full_image_slide(prs,
    "The Challenge",
    "10M+ Indians need palliative care | Only 1-2% have access\nLanguage barriers • Literacy gaps • 24/7 support needed")

# Slide 3: Grant & Funding
add_grant_slide(prs)

# Slide 4: Clinical Partners
add_partner_slide(prs)

# Slide 5: Video Demo - Marathi Interaction
add_video_slide(prs)

# Slide 6: Our Solution
add_content_slide(prs,
    "Solution: India's First Voice AI Palliative Care Helpline",
    [
        "Multi-modal access: Phone, WhatsApp, Web Voice",
        "7+ Indian languages: Hindi, Marathi, Bengali, Tamil, Punjabi, Malayalam, English",
        "3 Voice AI providers for redundancy & coverage",
        "RAG-powered knowledge base from Indian case vignettes",
        "Co-designed with Max Healthcare & Pallium India",
        "No PHI stored • HIPAA-aligned • Session expires in 24h",
    ])

# Slide 7: Architecture
add_architecture_slide(prs)

# Slide 8: Safety Features (NEW - today's work)
add_safety_slide(prs)

# Slide 9: Terminal Demo - Actual Code Running
add_demo_video_slide(prs)

# Slide 10: Demo Flow
add_demo_slide(prs)

# Slide 11: Impact
add_impact_slide(prs)

# Slide 12: Open Source
add_open_source_slide(prs)

# Slide 13: Thank You
add_thank_you_slide(prs)

# Save
output_path = "Palli_Sahayak_EkStep_V2_Final.pptx"
prs.save(output_path)

print(f"✅ Presentation created: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"\n🎯 KEY UPDATES:")
print(f"   • GCI Grant details (Dr Anurag Agrawal - PI, Ashish Makani - Co-I)")
print(f"   • Funders: GF India & BIRAC-DBT")
print(f"   • Clinical partners: Max Healthcare & Pallium India")
print(f"   • Marathi demo video reference")
print(f"   • Safety features (implemented today)")
print(f"   • Ashoka & KCDHA affiliation")
print(f"\n🚀 Ready for EkStep Voice AI Event!")
